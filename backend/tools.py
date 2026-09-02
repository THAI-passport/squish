"""Squish PDF tool implementations.

Every tool is a pure function over a temp working directory:

    fn(work: Path, inputs: list[Path], p: dict) -> Result

It writes its output inside `work` and returns a Result describing the file to
stream back. Nothing persists: app.py deletes `work` after the response is sent.

Heavy lifting is delegated to battle-tested engines rather than reimplemented:
  PyMuPDF      page surgery, rendering, redaction, overlays
  Ghostscript  compression, PDF/A
  LibreOffice  Office <-> PDF
  ocrmypdf     OCR (wraps tesseract + ghostscript)
  pdf2docx     PDF -> Word with layout reconstruction
"""

from __future__ import annotations

import glob
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable

import fitz  # PyMuPDF

log = logging.getLogger("uvicorn.error")

# ------------------------------------------------------------- budgets ---
# MAX_UPLOAD_MB caps what comes IN. These cap what goes OUT, which is the
# dimension that actually causes trouble: a 40,000-page PDF is small on disk
# but rasterises to tens of gigabytes, and /tmp is a tmpfs (RAM) in both the
# compose and Kubernetes deployments.
MAX_PAGES = int(os.environ.get("MAX_PAGES", "5000"))
# Total pixels a single job may render, in megapixels. 4000 MP is roughly
# 250 A4 pages at 150 DPI, or 1000 pages at 75 DPI after clamping.
MAX_RENDER_MP = float(os.environ.get("MAX_RENDER_MP", "4000"))
MIN_DPI = 36

PDF = "application/pdf"
ZIP = "application/zip"
DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class ToolError(Exception):
    """Raised for user-fixable problems; surfaced as HTTP 400."""


@dataclass
class Result:
    path: Path
    media_type: str
    filename: str


# ---------------------------------------------------------------- helpers ---

# Known alternative binary names on Windows
_WINDOWS_ENGINE_ALIASES: dict[str, list[str]] = {
    "gs": ["gswin64c", "gswin32c", "gs"],
    "soffice": ["soffice"],
    "tesseract": ["tesseract"],
    "ocrmypdf": ["ocrmypdf"],
    "qpdf": ["qpdf"],
    "exiftool": ["exiftool", "exiftool(-k)"],
}


def _windows_standard_paths(name: str) -> list[str]:
    """Known installation locations on Windows when binaries are not in PATH."""
    program_files = [
        os.environ.get("ProgramFiles", r"C:\Program Files"),
        os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)"),
    ]
    local_app_data = os.environ.get("LOCALAPPDATA")
    candidates: list[str] = []

    if name == "gs":
        for pf in program_files:
            candidates.extend(glob.glob(os.path.join(pf, "gs", "gs*", "bin", "gswin64c.exe")))
            candidates.extend(glob.glob(os.path.join(pf, "gs", "gs*", "bin", "gswin32c.exe")))
            candidates.extend(glob.glob(os.path.join(pf, "gs", "gs*", "bin", "gs.exe")))
    elif name == "soffice":
        for pf in program_files:
            candidates.append(os.path.join(pf, "LibreOffice", "program", "soffice.exe"))
            candidates.append(os.path.join(pf, "LibreOffice", "program", "soffice.com"))
    elif name == "tesseract":
        for pf in program_files:
            candidates.append(os.path.join(pf, "Tesseract-OCR", "tesseract.exe"))
        if local_app_data:
            candidates.append(os.path.join(local_app_data, "Programs", "Tesseract-OCR", "tesseract.exe"))
    elif name == "qpdf":
        for pf in program_files:
            candidates.append(os.path.join(pf, "qpdf", "bin", "qpdf.exe"))
    elif name == "ocrmypdf":
        # Check python environment Scripts directory
        candidates.append(os.path.join(sys.prefix, "Scripts", "ocrmypdf.exe"))
    elif name == "exiftool":
        for pf in program_files:
            candidates.append(os.path.join(pf, "exiftool", "exiftool.exe"))

    return [c for c in candidates if os.path.isfile(c)]


def find_engine(name: str, which: Callable[[str], str | None] | None = None) -> str | None:
    """Find the path to an engine binary, with Windows alias and path discovery.

    On Linux and macOS this delegates to `shutil.which(name)`. On Windows,
    Ghostscript is typically `gswin64c.exe`, LibreOffice lives in
    `C:\\Program Files\\LibreOffice\\program\\soffice.exe` without being added
    to PATH by default, and other tools have standard Windows locations.
    """
    _which = which or shutil.which

    # 1. Direct which lookup on PATH
    found = _which(name)
    if found:
        return found

    # If not on Windows, nothing more to check
    if os.name != "nt":
        return None

    # Map alias to canonical engine name if needed
    canonical = name
    for eng, aliases in _WINDOWS_ENGINE_ALIASES.items():
        if name in aliases or name == eng:
            canonical = eng
            break

    # 2. Check Windows aliases on PATH
    aliases = _WINDOWS_ENGINE_ALIASES.get(canonical, [])
    for alias in aliases:
        found = _which(alias)
        if found:
            return found

    # 3. Check known standard installation directories on Windows
    std_paths = _windows_standard_paths(canonical)
    if std_paths:
        return std_paths[0]

    return None


def run(cmd: list[str], cwd: Path | None = None, timeout: int = 300,
        ok_codes: tuple[int, ...] = (0,)) -> None:
    """Run a subprocess from an exec array (never a shell string).

    Subprocess stderr is captured into the uvicorn logger rather than
    discarded -- silent engine failures are undebuggable in production.

    `ok_codes` exists because qpdf exits 3 when it *successfully* recovered a
    damaged file. Treating that as failure would make the repair tool reject
    exactly the documents it is meant to fix.
    """
    binary = cmd[0]
    resolved = find_engine(binary)
    if resolved:
        cmd = [resolved] + cmd[1:]

    env = dict(os.environ)
    # LibreOffice and Ghostscript both need a writable HOME. Under a read-only
    # root filesystem on Kubernetes only /tmp is writable, so point HOME there.
    # On Windows, /tmp does not exist, so use tempfile.gettempdir() and ensure
    # USERPROFILE, TEMP, and TMP are also populated.
    default_home = str(cwd or tempfile.gettempdir())
    env.setdefault("HOME", default_home)
    if os.name == "nt":
        env.setdefault("USERPROFILE", default_home)
        env.setdefault("TEMP", tempfile.gettempdir())
        env.setdefault("TMP", tempfile.gettempdir())

    argv, preexec = _wrap_limits(cmd)
    log.info("exec: %s", " ".join(cmd[:4]) + (" ..." if len(cmd) > 4 else ""))
    try:
        proc = subprocess.run(
            argv, cwd=cwd, env=env, timeout=timeout,
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            preexec_fn=preexec,
        )
    except subprocess.TimeoutExpired:
        raise ToolError(f"{binary} timed out after {timeout}s")
    except FileNotFoundError:
        raise ToolError(f"{binary} is not installed in this image")
    if proc.returncode not in ok_codes:
        err = proc.stderr.decode("utf-8", "replace").strip()
        log.error("%s failed rc=%s: %s", binary, proc.returncode, err[-2000:])
        # Quote the shortest decisive line back to the user, not the whole dump.
        last = err.splitlines()[-1] if err else "no stderr"
        raise ToolError(f"{binary} failed: {last[:300]}")


# Ceilings applied to engine child processes. Ghostscript in particular has a
# long history of memory-exhaustion bugs on malformed input; a hostile PDF
# should kill its own subprocess, not the pod. Set to 0 to disable.
SUBPROC_MEM_MB = int(os.environ.get("SUBPROC_MEM_MB", "1536"))
SUBPROC_CPU_SEC = int(os.environ.get("SUBPROC_CPU_SEC", "600"))


def _wrap_limits(cmd: list[str]) -> tuple[list[str], Callable[[], None] | None]:
    """Apply the child resource ceilings by the safest means available.

    Tools run inside asyncio.to_thread, so the process forking these children
    is multi-threaded, and CPython documents preexec_fn as unsafe there: the
    child runs arbitrary Python between fork and exec while it may hold a lock
    another thread was using, which can deadlock the child forever.

    util-linux `prlimit` sets the same RLIMITs from outside, with no code
    running in the forked child at all -- so prefer it wherever it exists,
    which is every Linux image this ships in. macOS has no prlimit, and native
    `run-local.sh` is the case that matters there, so keep preexec_fn as the
    fallback: a small deadlock risk beats no memory ceiling at all.

    Windows has neither prlimit nor preexec_fn -- passing preexec_fn on Windows
    raises ValueError -- and the `resource` module does not exist there, so the
    limits are simply unavailable. Return None; the engine still runs, just
    without an in-process memory ceiling (Docker is the recommended Windows
    path precisely because the container can cap memory instead).
    """
    if SUBPROC_MEM_MB <= 0 and SUBPROC_CPU_SEC <= 0:
        return cmd, None
    prlimit = shutil.which("prlimit")
    if prlimit:
        wrapped = [prlimit]
        if SUBPROC_MEM_MB > 0:
            wrapped.append(f"--as={SUBPROC_MEM_MB * 1024 * 1024}")
        if SUBPROC_CPU_SEC > 0:
            wrapped.append(f"--cpu={SUBPROC_CPU_SEC}")
        wrapped.append("--core=0")
        # `--` keeps prlimit from eating flags that belong to the real command.
        return wrapped + ["--"] + cmd, None
    if os.name != "posix":          # Windows: no prlimit, no preexec_fn
        return cmd, None
    return cmd, _limits


def _limits() -> None:
    """Fallback path only (no prlimit). Applied in the child between fork and exec.

    RLIMIT_AS bounds address space and RLIMIT_CPU bounds CPU seconds, so a
    runaway engine dies on its own rather than being reaped by the OOM killer
    -- which on Kubernetes takes the whole container with it, including the
    other requests that pod was serving.
    """
    try:
        import resource
        if SUBPROC_MEM_MB > 0:
            b = SUBPROC_MEM_MB * 1024 * 1024
            resource.setrlimit(resource.RLIMIT_AS, (b, b))
        if SUBPROC_CPU_SEC > 0:
            resource.setrlimit(resource.RLIMIT_CPU,
                               (SUBPROC_CPU_SEC, SUBPROC_CPU_SEC))
        resource.setrlimit(resource.RLIMIT_CORE, (0, 0))
    except Exception:      # non-Linux, or limits already lower than requested
        pass


def open_pdf(path: Path, password: str = "") -> fitz.Document:
    try:
        doc = fitz.open(path)
    except Exception as exc:
        raise ToolError(f"cannot read {path.name}: {exc}")
    if doc.needs_pass:
        if not password or not doc.authenticate(password):
            doc.close()
            raise ToolError(f"{path.name} is password protected -- supply the password")
    if doc.page_count > MAX_PAGES:
        n = doc.page_count
        doc.close()
        raise ToolError(f"{path.name} has {n} pages; the limit is {MAX_PAGES}")
    return doc


def clamp_dpi(pages: list[int], doc: fitz.Document, dpi: int) -> int:
    """Reduce DPI until the whole job fits the render budget.

    Rendering is the one operation where a small input produces an unbounded
    output: an A4 page at 400 DPI is a 3308x4676 pixmap, about 46 MB
    uncompressed, and nothing about the input file size predicts the page
    count. Rather than refusing the job outright, scale the resolution down so
    it completes, and say so in the log.
    """
    if not pages:
        return dpi
    # Sum actual page areas: a document of A0 posters costs far more than A4.
    sq_inches = 0.0
    for i in pages:
        r = doc[i].rect
        sq_inches += (r.width / 72.0) * (r.height / 72.0)
    budget_px = MAX_RENDER_MP * 1_000_000
    wanted_px = sq_inches * dpi * dpi
    if wanted_px <= budget_px:
        return dpi
    safe = int((budget_px / sq_inches) ** 0.5)
    if safe < MIN_DPI:
        raise ToolError(
            f"rendering {len(pages)} pages exceeds the limit even at {MIN_DPI} DPI "
            f"-- select fewer pages"
        )
    log.warning("render budget: clamping %s DPI to %s for %d pages", dpi, safe, len(pages))
    return safe


def parse_pages(spec: str, total: int) -> list[int]:
    """'1-3,5,8-' -> zero-based page indices, deduped, in the order written.

    Empty spec means every page. Raises on out-of-range so the user gets a
    clear error instead of a silently truncated document.
    """
    spec = (spec or "").strip()
    if not spec:
        return list(range(total))
    out: list[int] = []
    for chunk in spec.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        m = re.fullmatch(r"(\d+)?\s*-\s*(\d+)?", chunk)
        if m:
            # A dash with a number on at least one side is a real range. A bare
            # "-" carries no bound and is almost always a typo, not "everything";
            # rejecting it beats silently selecting the whole document.
            if not m.group(1) and not m.group(2):
                raise ToolError(f"bad page range: {chunk!r}")
            start = int(m.group(1)) if m.group(1) else 1
            end = int(m.group(2)) if m.group(2) else total
        elif chunk.isdigit():
            start = end = int(chunk)
        else:
            raise ToolError(f"bad page range: {chunk!r}")
        if start < 1 or end > total or start > end:
            raise ToolError(f"page range {chunk!r} outside document (1-{total})")
        out.extend(range(start - 1, end))
    seen, uniq = set(), []
    for i in out:
        if i not in seen:
            seen.add(i)
            uniq.append(i)
    return uniq


def zip_dir(files: list[Path], dest: Path) -> Path:
    """Zip by basename, de-duplicating collisions.

    Two attachments both called `invoice.pdf` used to write the same archive
    name twice: most extractors keep only the last, so the user silently lost a
    file. Suffix repeats instead.
    """
    used: set[str] = set()
    with zipfile.ZipFile(dest, "w", zipfile.ZIP_DEFLATED) as z:
        for f in sorted(files):
            name = f.name
            if name in used:
                i = 2
                while f"{f.stem}({i}){f.suffix}" in used:
                    i += 1
                name = f"{f.stem}({i}){f.suffix}"
            used.add(name)
            z.write(f, name)
    return dest


def stem(path: Path) -> str:
    return re.sub(r"[^\w.\- ]+", "_", path.stem)[:80] or "document"


def save(doc: fitz.Document, out: Path, shrink: bool = True) -> Path:
    """Save with the flags that actually matter for size.

    garbage=4 runs full cross-reference garbage collection and merges
    duplicate objects; deflate recompresses streams. Without these, a merge of
    ten files that share a font embeds that font ten times.
    """
    doc.save(out, garbage=4 if shrink else 0, deflate=True, clean=shrink)
    return out


# ------------------------------------------------------------- organize ---

def merge_default_name(inputs: list[Path]) -> str:
    """A name that says what went in: 'a+b' for two, 'a+3-more' beyond that.

    Mirrors smartOutputName() in the UI so the field just previews the same
    value the server would pick when the field is left blank.
    """
    stems = [stem(f) for f in inputs]
    stems = [s for s in stems if s]
    if not stems:
        return "merged"
    if len(stems) == 1:
        return stems[0]
    if len(stems) == 2:
        return f"{stems[0]}+{stems[1]}"
    return f"{stems[0]}+{len(stems) - 1}-more"


def output_pdf_name(raw: str, fallback: str) -> str:
    """Sanitise a user-supplied output name into a safe `*.pdf` filename.

    The name reaches Content-Disposition and the temp path, so it must not
    escape the directory or carry control characters -- reuse safe_component,
    which already strips traversal, then guarantee the extension.
    """
    name = safe_component(raw) if raw and raw.strip() else fallback
    name = re.sub(r"\.pdf$", "", name, flags=re.IGNORECASE) or fallback
    return f"{name[:100]}.pdf"


def merge(work: Path, inputs: list[Path], p: dict) -> Result:
    if len(inputs) < 2:
        raise ToolError("merge needs at least 2 files")
    out = fitz.open()
    for f in inputs:
        src = open_pdf(f, p.get("password", ""))
        out.insert_pdf(src)
        src.close()
    name = output_pdf_name(p.get("output_name", ""), merge_default_name(inputs))
    dest = save(out, work / name)
    out.close()
    return Result(dest, PDF, name)


def split(work: Path, inputs: list[Path], p: dict) -> Result:
    """mode=ranges  -> one PDF containing the selected pages
       mode=every   -> one PDF per page, zipped
       mode=chunks  -> fixed-size groups of pages, zipped
    """
    src = open_pdf(inputs[0], p.get("password", ""))
    total = src.page_count
    base = stem(inputs[0])
    mode = p.get("mode", "ranges")

    if mode == "ranges":
        pages = parse_pages(p.get("pages", ""), total)
        out = fitz.open()
        for i in pages:
            out.insert_pdf(src, from_page=i, to_page=i)
        dest = save(out, work / f"{base}_pages.pdf")
        out.close()
        src.close()
        return Result(dest, PDF, f"{base}_pages.pdf")

    size = 1 if mode == "every" else max(1, int(p.get("size", 1) or 1))
    parts: list[Path] = []
    for start in range(0, total, size):
        end = min(start + size - 1, total - 1)
        out = fitz.open()
        out.insert_pdf(src, from_page=start, to_page=end)
        label = f"{start + 1}" if start == end else f"{start + 1}-{end + 1}"
        part = save(out, work / f"{base}_{label}.pdf")
        out.close()
        parts.append(part)
    src.close()
    dest = zip_dir(parts, work / f"{base}_split.zip")
    return Result(dest, ZIP, f"{base}_split.zip")


def remove_pages(work: Path, inputs: list[Path], p: dict) -> Result:
    src = open_pdf(inputs[0], p.get("password", ""))
    # A blank spec means "every page" everywhere else, which here would mean
    # deleting the whole document. Ask instead of reporting it as an error
    # about removing every page, which reads like a bug.
    if not (p.get("pages") or "").strip():
        src.close()
        raise ToolError("name the pages to remove, e.g. 2 or 4-6")
    drop = set(parse_pages(p.get("pages", ""), src.page_count))
    if not drop:
        raise ToolError("select at least one page to remove")
    keep = [i for i in range(src.page_count) if i not in drop]
    if not keep:
        raise ToolError("that would remove every page")
    out = fitz.open()
    for i in keep:
        out.insert_pdf(src, from_page=i, to_page=i)
    base = stem(inputs[0])
    dest = save(out, work / f"{base}_trimmed.pdf")
    out.close()
    src.close()
    return Result(dest, PDF, f"{base}_trimmed.pdf")


def organize(work: Path, inputs: list[Path], p: dict) -> Result:
    """Reorder pages to an explicit sequence, e.g. '3,1,2,4-'."""
    src = open_pdf(inputs[0], p.get("password", ""))
    order = parse_pages(p.get("pages", ""), src.page_count)
    out = fitz.open()
    for i in order:
        out.insert_pdf(src, from_page=i, to_page=i)
    base = stem(inputs[0])
    dest = save(out, work / f"{base}_organized.pdf")
    out.close()
    src.close()
    return Result(dest, PDF, f"{base}_organized.pdf")


def rotate(work: Path, inputs: list[Path], p: dict) -> Result:
    src = open_pdf(inputs[0], p.get("password", ""))
    angle = int(p.get("angle", 90))
    if angle % 90:
        raise ToolError("angle must be a multiple of 90")
    targets = set(parse_pages(p.get("pages", ""), src.page_count))
    for i in targets:
        page = src[i]
        page.set_rotation((page.rotation + angle) % 360)
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_rotated.pdf")
    src.close()
    return Result(dest, PDF, f"{base}_rotated.pdf")


# ------------------------------------------------------------- optimize ---

# Ghostscript's presets, cheapest to best. /screen downsamples images to 72dpi,
# /ebook to 150dpi, /printer to 300dpi.
GS_LEVELS = {"extreme": "/screen", "recommended": "/ebook", "low": "/printer"}


def compress(work: Path, inputs: list[Path], p: dict) -> Result:
    if not inputs[0].stat().st_size:
        raise ToolError(f"cannot read {inputs[0].name}: file is empty")
    level = GS_LEVELS.get(p.get("level", "recommended"), "/ebook")
    base = stem(inputs[0])
    dest = work / f"{base}_compressed.pdf"
    run([
        "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.7",
        f"-dPDFSETTINGS={level}", "-dNOPAUSE", "-dQUIET", "-dBATCH",
        "-dDetectDuplicateImages=true", "-dCompressFonts=true",
        f"-sOutputFile={dest}", str(inputs[0]),
    ], cwd=work)
    # Ghostscript can inflate an already-optimised file. Keep the smaller one
    # rather than handing back a "compressed" file that is bigger.
    if dest.stat().st_size >= inputs[0].stat().st_size:
        log.info("gs output larger than input; returning original")
        shutil.copy(inputs[0], dest)
    return Result(dest, PDF, f"{base}_compressed.pdf")


def compress_wasm(work: Path, inputs: list[Path], p: dict) -> Result:
    """Browser-safe structural compression using PyMuPDF only.

    WebAssembly cannot run Ghostscript, so static deployments cannot perform
    lossy image downsampling. They can still garbage-collect unreachable
    objects, deduplicate streams, compress streams, and generate object
    streams. As with the server implementation, never return a larger file.
    """
    source = inputs[0]
    if not source.stat().st_size:
        raise ToolError(f"cannot read {source.name}: file is empty")
    base = stem(source)
    dest = work / f"{base}_compressed.pdf"
    src = open_pdf(source, p.get("password", ""))
    try:
        src.save(dest, garbage=4, deflate=True, clean=True, use_objstms=1)
    finally:
        src.close()
    if dest.stat().st_size >= source.stat().st_size:
        shutil.copyfile(source, dest)
    return Result(dest, PDF, f"{base}_compressed.pdf")


def repair(work: Path, inputs: list[Path], p: dict) -> Result:
    """Rebuild a damaged cross-reference table.

    qpdf reconstructs the xref by scanning for object headers, recovering far
    more than a plain reparse. Two details matter:
      * rc 3 means "recovered, with warnings" -- the success case here.
      * do NOT pass --qdf: that emits the uncompressed debugging form, which
        roughly doubles file size for no benefit to the user.
    """
    base = stem(inputs[0])
    dest = work / f"{base}_repaired.pdf"
    try:
        run(["qpdf", "--object-streams=generate", "--stream-data=compress",
             str(inputs[0]), str(dest)], cwd=work, ok_codes=(0, 3))
    except ToolError as exc:
        log.warning("qpdf repair failed (%s), falling back to PyMuPDF", exc)
        src = fitz.open(inputs[0])   # fitz repairs on open where it can
        save(src, dest)
        src.close()
    if not dest.exists() or dest.stat().st_size == 0:
        raise ToolError("the file is too damaged to recover")
    return Result(dest, PDF, f"{base}_repaired.pdf")


def repair_wasm(work: Path, inputs: list[Path], p: dict) -> Result:
    """Browser-safe best-effort xref rebuild.

    PyMuPDF repairs many broken cross-reference tables while opening a PDF.
    qpdf remains the stronger server engine, but this recovery pass is useful
    and honest for Cloudflare's browser-only build.
    """
    base = stem(inputs[0])
    dest = work / f"{base}_repaired.pdf"
    try:
        src = fitz.open(inputs[0])
        if src.needs_pass:
            password = str(p.get("password") or "")
            if not password or not src.authenticate(password):
                src.close()
                raise ToolError("this PDF is password protected")
        src.save(dest, garbage=4, deflate=True, clean=True, use_objstms=1)
        src.close()
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"the file is too damaged to recover: {exc}")
    if not dest.exists() or dest.stat().st_size == 0:
        raise ToolError("the file is too damaged to recover")
    return Result(dest, PDF, f"{base}_repaired.pdf")


def ocr(work: Path, inputs: list[Path], p: dict) -> Result:
    """Make a scanned PDF searchable. Text layer is added, image is kept."""
    base = stem(inputs[0])
    dest = work / f"{base}_ocr.pdf"
    lang = re.sub(r"[^a-z+]", "", str(p.get("lang", "eng")).lower()) or "eng"
    cmd = ["ocrmypdf", "--output-type", "pdf", "-l", lang, "--jobs", "2"]
    if p.get("force"):
        # Re-OCR pages that already carry a text layer.
        cmd.append("--force-ocr")
    else:
        # Default: skip pages that already have text instead of erroring out.
        cmd.append("--skip-text")
    if p.get("deskew"):
        cmd.append("--deskew")
    run(cmd + [str(inputs[0]), str(dest)], cwd=work, timeout=900)
    return Result(dest, PDF, f"{base}_ocr.pdf")


# -------------------------------------------------------------- convert ---

def pdf_to_jpg(work: Path, inputs: list[Path], p: dict) -> Result:
    src = open_pdf(inputs[0], p.get("password", ""))
    dpi = max(MIN_DPI, min(600, int(p.get("dpi", 150))))
    fmt = "png" if p.get("format") == "png" else "jpg"
    base = stem(inputs[0])
    pages = parse_pages(p.get("pages", ""), src.page_count)
    dpi = clamp_dpi(pages, src, dpi)
    # jpg_quality must be omitted for PNG, not passed as None: PyMuPDF hands the
    # value straight to the C layer, which wants an int.
    opts: dict[str, Any] = {}
    if fmt == "jpg":
        opts["jpg_quality"] = max(1, min(100, int(p.get("quality", 85))))
    files: list[Path] = []
    for i in pages:
        pix = src[i].get_pixmap(dpi=dpi)
        f = work / f"{base}_{i + 1:04d}.{fmt}"
        pix.save(f, **opts)
        files.append(f)
    src.close()
    if not files:
        raise ToolError("no pages selected")
    if len(files) == 1:
        mt = "image/png" if fmt == "png" else "image/jpeg"
        return Result(files[0], mt, files[0].name)
    dest = zip_dir(files, work / f"{base}_images.zip")
    return Result(dest, ZIP, f"{base}_images.zip")


def jpg_to_pdf(work: Path, inputs: list[Path], p: dict) -> Result:
    """One image per page. Images are fitted inside the page, not stretched."""
    sizes = {
        "fit": None,
        "a4": fitz.paper_rect("a4"),
        "letter": fitz.paper_rect("letter"),
    }
    target = sizes.get(p.get("size", "fit"), None)
    margin = max(0, min(200, int(p.get("margin", 0))))
    out = fitz.open()
    for f in inputs:
        try:
            img = fitz.open(f)
            rect = img[0].rect
            pdf_bytes = img.convert_to_pdf()
            img.close()
        except Exception as exc:
            raise ToolError(f"{f.name} is not a readable image: {exc}")
        src = fitz.open("pdf", pdf_bytes)
        if target is None:
            page = out.new_page(width=rect.width, height=rect.height)
            box = page.rect
        else:
            page = out.new_page(width=target.width, height=target.height)
            box = page.rect + (margin, margin, -margin, -margin)
            # Preserve aspect ratio inside the margin box.
            scale = min(box.width / rect.width, box.height / rect.height)
            w, h = rect.width * scale, rect.height * scale
            cx, cy = box.x0 + box.width / 2, box.y0 + box.height / 2
            box = fitz.Rect(cx - w / 2, cy - h / 2, cx + w / 2, cy + h / 2)
        page.show_pdf_page(box, src, 0)
        src.close()
    if not out.page_count:
        raise ToolError("no images supplied")
    dest = save(out, work / "images.pdf")
    out.close()
    return Result(dest, PDF, "images.pdf")


OFFICE_EXT = {".doc", ".docx", ".odt", ".rtf", ".txt",
              ".xls", ".xlsx", ".ods", ".csv",
              ".ppt", ".pptx", ".odp"}


LO_TEMPLATE = Path(os.environ.get("LO_PROFILE_TEMPLATE", "/opt/lo-profile"))


def office_to_pdf(work: Path, inputs: list[Path], p: dict) -> Result:
    """Word/Excel/PowerPoint -> PDF via headless LibreOffice.

    Each request gets its OWN profile directory, because concurrent soffice
    processes sharing one profile corrupt it. Building a profile from scratch
    costs 3-5 seconds, so the image bakes a warm one at /opt/lo-profile and we
    copy that instead -- a few MB of file copy versus a full LibreOffice
    first-run. Falls back to cold creation if the template is absent (native
    `run-local.sh`, or a custom image).
    """
    outs: list[Path] = []
    profile = work / "lo_profile"
    if LO_TEMPLATE.is_dir():
        try:
            shutil.copytree(LO_TEMPLATE, profile)
        except Exception as exc:
            log.warning("could not seed LibreOffice profile: %s", exc)
    for f in inputs:
        if f.suffix.lower() not in OFFICE_EXT:
            raise ToolError(f"{f.name}: unsupported type {f.suffix}")
        run([
            "soffice", "--headless", "--norestore", "--nolockcheck",
            f"-env:UserInstallation={profile.as_uri()}",
            "--convert-to", "pdf", "--outdir", str(work), str(f),
        ], cwd=work, timeout=600)
        produced = work / (f.stem + ".pdf")
        if not produced.exists():
            raise ToolError(f"LibreOffice produced no output for {f.name}")
        outs.append(produced)
    if len(outs) == 1:
        return Result(outs[0], PDF, outs[0].name)
    dest = zip_dir(outs, work / "converted.zip")
    return Result(dest, ZIP, "converted.zip")


# Markdown extensions chosen for fidelity: GitHub-style tables, fenced code
# with syntax highlighting, footnotes, definition lists, task lists, a TOC
# anchor, and raw HTML passthrough. `extra` already bundles tables, fenced_code,
# footnotes, attr_list, def_list and abbr.
MD_EXTENSIONS = [
    "extra", "codehilite", "sane_lists", "toc", "admonition",
]

# GitHub-flavoured print stylesheet. The print-specific rules are what make the
# output faithful rather than merely styled:
#   * @page sets a real paper size and margins
#   * thead { display: table-header-group } repeats table headers across page
#     breaks -- the single most common "lost table" failure
#   * tr / pre / figure avoid being split mid-element
#   * pre wraps instead of overflowing the page width
MD_CSS = """
@page { size: %(page)s; margin: %(margin)smm; }
* { box-sizing: border-box; }
body { font-family: -apple-system, "Segoe UI", Roboto, "Helvetica Neue", Arial,
       "DejaVu Sans", sans-serif; font-size: 11pt; line-height: 1.55;
       color: #1f2328; }
h1,h2,h3,h4,h5,h6 { font-weight: 600; line-height: 1.25; margin: 1.2em 0 .5em;
       break-after: avoid; }
h1 { font-size: 2em; border-bottom: 1px solid #d1d9e0; padding-bottom: .3em; }
h2 { font-size: 1.5em; border-bottom: 1px solid #d1d9e0; padding-bottom: .3em; }
h3 { font-size: 1.25em; } h4 { font-size: 1em; }
p, ul, ol, blockquote, table, pre { margin: 0 0 1em; }
a { color: #0969da; text-decoration: none; }
code { font-family: "DejaVu Sans Mono", ui-monospace, SFMono-Regular, Menlo,
       Consolas, monospace; font-size: .88em;
       background: #eff1f3; padding: .15em .35em; border-radius: 4px; }
pre { background: #f6f8fa; padding: 12px 14px; border-radius: 6px;
      overflow-wrap: anywhere; white-space: pre-wrap; break-inside: avoid; }
pre code { background: none; padding: 0; }
blockquote { color: #59636e; border-left: .25em solid #d1d9e0; padding: 0 1em;
      margin-left: 0; }
img { max-width: 100%%; }
figure, tr { break-inside: avoid; }
table { border-collapse: collapse; width: 100%%; font-size: .95em; }
thead { display: table-header-group; }
th, td { border: 1px solid #d1d9e0; padding: 6px 13px; text-align: left;
      vertical-align: top; }
th { background: #f6f8fa; font-weight: 600; }
tbody tr:nth-child(even) { background: #f6f8fa; }
hr { border: 0; border-top: 1px solid #d1d9e0; margin: 1.5em 0; }
"""


def _weasy_url_fetcher(work: Path, allow_remote: bool):
    """Restrict what a Markdown document can pull in when rendered.

    WeasyPrint would otherwise happily fetch `file:///etc/passwd` or
    `http://169.254.169.254/` (cloud metadata) referenced by a hostile
    document. So: allow `data:` URIs always; block local files entirely (a
    single-file upload has no legitimate local assets); and permit http(s) only
    when the user opts in, and even then reject private/loopback/link-local
    addresses to blunt SSRF. This is imperfect against DNS-rebinding -- the
    default posture (remote OFF) is the safe one.
    """
    import ipaddress
    import socket
    from urllib.parse import urlparse
    from weasyprint import default_url_fetcher

    def fetch(url: str, *args, **kwargs):
        if url.startswith("data:"):
            return default_url_fetcher(url, *args, **kwargs)
        scheme = urlparse(url).scheme
        if scheme in ("http", "https"):
            if not allow_remote:
                raise ValueError("remote resources are disabled for this document")
            host = urlparse(url).hostname or ""
            for info in socket.getaddrinfo(host, None):
                ip = ipaddress.ip_address(info[4][0])
                if (ip.is_private or ip.is_loopback or ip.is_link_local
                        or ip.is_reserved or ip.is_multicast):
                    raise ValueError(f"blocked non-public address for {host}")
            return default_url_fetcher(url, *args, **kwargs)
        # file:// and anything else: refused, so a document cannot read the
        # server's filesystem.
        raise ValueError(f"blocked resource scheme: {scheme or 'file'}")

    return fetch


def md_to_pdf(work: Path, inputs: list[Path], p: dict) -> Result:
    """Render a Markdown file to PDF via HTML + WeasyPrint.

    Fidelity comes from a real CSS print engine, not a lossy converter: tables
    keep their borders and repeat their header row across page breaks, fenced
    code is syntax-highlighted and wraps instead of overflowing, and page size
    and margins are honoured.
    """
    import markdown
    from pygments.formatters import HtmlFormatter
    try:
        from weasyprint import HTML
        import weasyprint
        has_weasy = True
    except ImportError:
        has_weasy = False

    # Input is either an uploaded file (wins if present) or pasted Markdown from
    # the text box. The tool is registered with min_files=0, so app.py allows a
    # request with no upload -- validate our own input here.
    if inputs:
        src = inputs[0]
        try:
            text = src.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = src.read_text(encoding="utf-8", errors="replace")
        base = stem(src)
    else:
        text = str(p.get("md_text") or "")
        if not text.strip():
            raise ToolError("upload a .md file or paste some Markdown to convert")
        base = safe_component(p.get("title") or "document")

    page = "Letter" if str(p.get("page_size", "a4")).lower() == "letter" else "A4"
    try:
        margin = max(0, min(50, int(float(p.get("margin", 18)))))
    except (TypeError, ValueError):
        margin = 18
    allow_remote = bool(p.get("allow_remote"))

    body_html = markdown.markdown(
        text, extensions=MD_EXTENSIONS, output_format="html5")
    # Pygments emits its own classes under .codehilite; pull the matching colours.
    pygments_css = HtmlFormatter().get_style_defs(".codehilite")
    css = (MD_CSS % {"page": page, "margin": margin}) + "\n" + pygments_css

    title = safe_component(p.get("title") or base)
    doc_html = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<title>{title}</title><style>{css}</style></head>"
        f"<body>{body_html}</body></html>"
    )

    dest = work / f"{base}.pdf"
    
    if has_weasy:
        HTML(string=body_html, base_url=str(work), url_fetcher=_weasy_url_fetcher(work, allow_remote)).write_pdf(
            target=str(dest), stylesheets=[weasyprint.CSS(string=css)],
            optimize_images=True, jpeg_quality=85, dpi=300,
            presentational_hints=True, optimize_size=("fonts", "images"))
    else:
        # Fallback to PyMuPDF Story for WebAssembly where WeasyPrint is unavailable.
        # This provides a basic render, albeit without full CSS page layout.
        import fitz
        html_doc = f"<html><head><style>{css}</style></head><body>{body_html}</body></html>"
        story = fitz.Story(html=html_doc)
        writer = fitz.DocumentWriter(str(dest))
        
        # A4 page size and rects
        page_rect = fitz.Rect(0, 0, 595, 842) if page == "A4" else fitz.Rect(0, 0, 612, 792)
        m_pt = margin * 2.83465  # mm to pt
        draw_rect = page_rect + (m_pt, m_pt, -m_pt, -m_pt)
        
        more = 1
        while more:
            device = writer.begin_page(page_rect)
            more, _ = story.place(draw_rect)
            story.draw(device)
            writer.end_page()
        writer.close()

    return Result(dest, PDF, f"{base}.pdf")


def md_to_pdf_wasm(work: Path, inputs: list[Path], p: dict) -> Result:
    """Browser-only Markdown renderer using PyMuPDF Story.

    Kept as a separate entry point so static-mode dependency checks can prove
    this path never imports WeasyPrint or its native Cairo/Pango stack.
    """
    import markdown
    from pygments.formatters import HtmlFormatter

    if inputs:
        src = inputs[0]
        try:
            text = src.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = src.read_text(encoding="utf-8", errors="replace")
        base = stem(src)
    else:
        text = str(p.get("md_text") or "")
        if not text.strip():
            raise ToolError("upload a .md file or paste some Markdown to convert")
        base = safe_component(p.get("title") or "document")

    page_name = "Letter" if str(p.get("page_size", "a4")).lower() == "letter" else "A4"
    try:
        margin = max(0, min(50, int(float(p.get("margin", 18)))))
    except (TypeError, ValueError):
        margin = 18
    body_html = markdown.markdown(text, extensions=MD_EXTENSIONS, output_format="html5")
    css = (MD_CSS % {"page": page_name, "margin": margin}) + "\n" + HtmlFormatter().get_style_defs(".codehilite")
    story = fitz.Story(html=f"<html><head><style>{css}</style></head><body>{body_html}</body></html>")
    dest = work / f"{base}.pdf"
    writer = fitz.DocumentWriter(str(dest))
    page_rect = fitz.Rect(0, 0, 595, 842) if page_name == "A4" else fitz.Rect(0, 0, 612, 792)
    m_pt = margin * 2.83465
    draw_rect = page_rect + (m_pt, m_pt, -m_pt, -m_pt)
    more = 1
    try:
        while more:
            device = writer.begin_page(page_rect)
            more, _ = story.place(draw_rect)
            story.draw(device)
            writer.end_page()
    finally:
        writer.close()
    return Result(dest, PDF, f"{base}.pdf")


def pdf_to_word(work: Path, inputs: list[Path], p: dict) -> Result:
    """Reconstructs paragraphs, tables and images. Complex layouts drift."""
    from pdf2docx import Converter
    base = stem(inputs[0])
    dest = work / f"{base}.docx"
    conv = Converter(str(inputs[0]), password=p.get("password") or None)
    try:
        conv.convert(str(dest), start=0, end=None)
    except Exception as exc:
        raise ToolError(f"conversion failed: {exc}")
    finally:
        conv.close()
    return Result(dest, DOCX, f"{base}.docx")


def pdf_to_excel(work: Path, inputs: list[Path], p: dict) -> Result:
    """Extracts detected tables, one worksheet per table.

    This is honest table extraction, not OCR: a PDF with no ruled or
    whitespace-aligned table structure will yield nothing.
    """
    from openpyxl import Workbook
    src = open_pdf(inputs[0], p.get("password", ""))
    wb = Workbook()
    wb.remove(wb.active)
    found = 0
    for pno in parse_pages(p.get("pages", ""), src.page_count):
        for t, table in enumerate(src[pno].find_tables().tables, 1):
            found += 1
            ws = wb.create_sheet(f"p{pno + 1}_t{t}"[:31])
            for row in table.extract():
                ws.append([("" if c is None else str(c)) for c in row])
    src.close()
    if not found:
        raise ToolError("no tables detected -- run OCR first if this is a scan")
    base = stem(inputs[0])
    dest = work / f"{base}.xlsx"
    wb.save(dest)
    return Result(dest, XLSX, f"{base}.xlsx")


def pdf_to_powerpoint(work: Path, inputs: list[Path], p: dict) -> Result:
    """One slide per page, rendered at full bleed. Faithful, not editable."""
    from pptx import Presentation
    from pptx.util import Emu
    src = open_pdf(inputs[0], p.get("password", ""))
    dpi = max(72, min(300, int(p.get("dpi", 150))))
    pages = parse_pages(p.get("pages", ""), src.page_count)
    dpi = clamp_dpi(pages, src, dpi)
    if not pages:
        src.close()
        raise ToolError("no pages selected")
    prs = Presentation()
    # Size the deck from the FIRST SELECTED page, not page 0 -- page 0 may not
    # even be in the selection. A deck has one slide size, so mixed page sizes
    # still have to be letterboxed rather than stretched.
    first = src[pages[0]].rect
    # PowerPoint measures in EMU; 914400 EMU per inch, PDF points are 1/72 inch.
    prs.slide_width = Emu(int(first.width / 72 * 914400))
    prs.slide_height = Emu(int(first.height / 72 * 914400))
    blank = prs.slide_layouts[6]
    for pno in pages:
        img = work / f"slide_{pno:04d}.png"
        src[pno].get_pixmap(dpi=dpi).save(img)
        slide = prs.slides.add_slide(blank)
        r = src[pno].rect
        # Fit inside the slide, preserving aspect ratio, centred.
        scale = min(prs.slide_width / (r.width / 72 * 914400),
                    prs.slide_height / (r.height / 72 * 914400))
        w = int(r.width / 72 * 914400 * scale)
        h = int(r.height / 72 * 914400 * scale)
        slide.shapes.add_picture(str(img), Emu(int((prs.slide_width - w) / 2)),
                                 Emu(int((prs.slide_height - h) / 2)),
                                 width=Emu(w), height=Emu(h))
    src.close()
    base = stem(inputs[0])
    dest = work / f"{base}.pptx"
    prs.save(dest)
    return Result(dest, PPTX, f"{base}.pptx")


def pdf_to_pdfa(work: Path, inputs: list[Path], p: dict) -> Result:
    """Archival PDF/A-2b. ocrmypdf drives Ghostscript's PDF/A pipeline."""
    base = stem(inputs[0])
    dest = work / f"{base}_pdfa.pdf"
    run(["ocrmypdf", "--output-type", "pdfa-2", "--skip-text",
         "--tesseract-timeout", "0", str(inputs[0]), str(dest)],
        cwd=work, timeout=900)
    return Result(dest, PDF, f"{base}_pdfa.pdf")


def pdf_to_markdown(work: Path, inputs: list[Path], p: dict) -> Result:
    src = open_pdf(inputs[0], p.get("password", ""))
    lines: list[str] = []
    for pno in parse_pages(p.get("pages", ""), src.page_count):
        lines.append(f"\n---\n\n<!-- page {pno + 1} -->\n")
        lines.append(src[pno].get_text("text").strip())
    src.close()
    base = stem(inputs[0])
    dest = work / f"{base}.md"
    dest.write_text("\n".join(lines).strip() + "\n", encoding="utf-8")
    return Result(dest, "text/markdown", f"{base}.md")


def pdf_to_text(work: Path, inputs: list[Path], p: dict) -> Result:
    src = open_pdf(inputs[0], p.get("password", ""))
    lines: list[str] = []
    for pno in parse_pages(p.get("pages", ""), src.page_count):
        lines.append(src[pno].get_text("text").strip())
    src.close()
    base = stem(inputs[0])
    dest = work / f"{base}.txt"
    dest.write_text("\n\n".join(lines).strip() + "\n", encoding="utf-8")
    return Result(dest, "text/plain", f"{base}.txt")


# ------------------------------------------------------------------ edit ---

POSITIONS = {
    "top-left": (0.10, 0.08), "top": (0.50, 0.08), "top-right": (0.90, 0.08),
    "left": (0.10, 0.50), "center": (0.50, 0.50), "right": (0.90, 0.50),
    "bottom-left": (0.10, 0.94), "bottom": (0.50, 0.94),
    "bottom-right": (0.90, 0.94),
}


def watermark(work: Path, inputs: list[Path], p: dict) -> Result:
    src = open_pdf(inputs[0], p.get("password", ""))
    text = str(p.get("text") or "").strip()
    if not text:
        raise ToolError("watermark text is required")
    size = max(6, min(200, int(p.get("size", 42))))
    opacity = max(0.05, min(1.0, float(p.get("opacity", 0.25))))
    angle = int(p.get("angle", 45))
    color = _hex_rgb(p.get("color") or "#7c5cff")
    fx, fy = POSITIONS.get(p.get("position", "center"), (0.5, 0.5))
    mode = p.get("mode", "single")  # single | tile

    for pno in parse_pages(p.get("pages", ""), src.page_count):
        page = src[pno]
        r = page.rect
        spots = []
        if mode == "tile":
            step_x, step_y = r.width / 3, r.height / 4
            for gx in range(3):
                for gy in range(4):
                    spots.append((step_x * (gx + 0.5), step_y * (gy + 0.5)))
        else:
            spots.append((r.width * fx, r.height * fy))
        for cx, cy in spots:
            # Morph rotates the text about its own insertion point; without the
            # pivot the glyphs rotate around the page origin and fly off-page.
            pivot = fitz.Point(cx, cy)
            morph = (pivot, fitz.Matrix(1, 1).prerotate(angle))
            page.insert_text(
                pivot - fitz.Point(len(text) * size * 0.25, 0),
                text, fontsize=size, fontname="helv", color=color,
                fill_opacity=opacity, stroke_opacity=opacity, morph=morph,
                overlay=bool(p.get("overlay", True)),
            )
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_watermarked.pdf")
    src.close()
    return Result(dest, PDF, f"{base}_watermarked.pdf")


def page_numbers(work: Path, inputs: list[Path], p: dict) -> Result:
    src = open_pdf(inputs[0], p.get("password", ""))
    start = int(p.get("start", 1))
    size = max(6, min(72, int(p.get("size", 11))))
    fmt = str(p.get("format", "{n}"))
    if "{n}" not in fmt:
        raise ToolError("format must contain {n}")
    fx, fy = POSITIONS.get(p.get("position", "bottom"), (0.5, 0.94))
    color = _hex_rgb(p.get("color") or "#000000")
    pages = parse_pages(p.get("pages", ""), src.page_count)
    # {total} is the document's page count, not the count of pages being
    # stamped -- numbering pages 1-3 of a 10-page file must not say "of 3".
    total = src.page_count
    for seq, pno in enumerate(pages):
        page = src[pno]
        label = fmt.replace("{n}", str(start + seq)).replace("{total}", str(total))
        r = page.rect
        w = fitz.get_text_length(label, fontname="helv", fontsize=size)
        page.insert_text(
            fitz.Point(r.width * fx - w / 2, r.height * fy),
            label, fontsize=size, fontname="helv", color=color, overlay=True,
        )
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_numbered.pdf")
    src.close()
    return Result(dest, PDF, f"{base}_numbered.pdf")


# set_cropbox works in UNROTATED page space, but the user picks margins against
# what they see on screen. On a /Rotate 90 page the visible top edge is the
# unrotated LEFT edge, so applying the margins verbatim trims the wrong sides.
# Map visible edge -> unrotated edge per rotation.
_ROT_EDGES = {
    0:   {"top": "top", "right": "right", "bottom": "bottom", "left": "left"},
    90:  {"top": "left", "right": "top", "bottom": "right", "left": "bottom"},
    180: {"top": "bottom", "right": "left", "bottom": "top", "left": "right"},
    270: {"top": "right", "right": "bottom", "bottom": "left", "left": "top"},
}


def crop(work: Path, inputs: list[Path], p: dict) -> Result:
    """Trim by margin in points from each edge, as the page is displayed."""
    src = open_pdf(inputs[0], p.get("password", ""))
    m = {k: float(p.get(k, 0) or 0) for k in ("top", "bottom", "left", "right")}
    for pno in parse_pages(p.get("pages", ""), src.page_count):
        page = src[pno]
        rot = page.rotation % 360
        mapping = _ROT_EDGES.get(rot)
        if mapping is None:                     # non-multiple-of-90 rotation
            raise ToolError(f"page {pno + 1} has an unsupported rotation ({rot})")
        # Un-rotate to get the coordinate space set_cropbox expects, then put
        # the rotation back afterwards so the page still displays as before.
        page.set_rotation(0)
        r = page.rect
        um = {unrot: m[vis] for vis, unrot in mapping.items()}
        box = fitz.Rect(r.x0 + um["left"], r.y0 + um["top"],
                        r.x1 - um["right"], r.y1 - um["bottom"])
        if box.width <= 1 or box.height <= 1:
            page.set_rotation(rot)
            raise ToolError(f"crop leaves nothing on page {pno + 1}")
        page.set_cropbox(box)
        page.set_rotation(rot)
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_cropped.pdf")
    src.close()
    return Result(dest, PDF, f"{base}_cropped.pdf")


# -------------------------------------------------------------- security ---

def protect(work: Path, inputs: list[Path], p: dict) -> Result:
    pw = str(p.get("password_new") or "")
    if len(pw) < 4:
        raise ToolError("password must be at least 4 characters")
    src = open_pdf(inputs[0], p.get("password", ""))
    perm = (fitz.PDF_PERM_ACCESSIBILITY | fitz.PDF_PERM_PRINT)
    if p.get("allow_copy"):
        perm |= fitz.PDF_PERM_COPY
    if p.get("allow_modify"):
        perm |= fitz.PDF_PERM_MODIFY | fitz.PDF_PERM_ANNOTATE
    base = stem(inputs[0])
    dest = work / f"{base}_protected.pdf"
    src.save(dest, encryption=fitz.PDF_ENCRYPT_AES_256,
             owner_pw=str(p.get("owner_password") or pw),
             user_pw=pw, permissions=perm, garbage=4, deflate=True)
    src.close()
    return Result(dest, PDF, f"{base}_protected.pdf")


def _dispatch_extract_pages(work: Path, source: Path, spec: str,
                            open_password: str = "") -> Path:
    """Create a page-only tree for one burst recipient and scrub catalog data."""
    src = open_pdf(source, open_password)
    pages = parse_pages(spec, src.page_count)
    if not pages:
        src.close()
        raise ToolError("dispatch page mapping cannot be empty")
    out = fitz.open()
    for page_no in pages:
        # Dispatch isolation deliberately removes links, annotations and form
        # widgets. They can retain document-level references to excluded pages.
        out.insert_pdf(src, from_page=page_no, to_page=page_no,
                       links=0, annots=0, widgets=0)
    src.close()
    out.scrub(
        attached_files=True, embedded_files=True, javascript=True,
        metadata=True, xml_metadata=True, remove_links=True,
        reset_fields=True, reset_responses=True, thumbnails=True,
        clean_pages=True, hidden_text=False, redactions=False,
    )
    dest = work / f"{stem(source)}_pages_{safe_component(spec)[:40]}.pdf"
    save(out, dest)
    out.close()
    return dest


def _encode_zero_width_tag(recipient: str, secret_seed: str) -> tuple[str, str]:
    """Encode a keyed leak marker into invisible zero-width characters.
    
    This is a correlation marker, not forensic proof: it can be stripped or
    copied. The caller must supply a deployment or per-dispatch secret so the
    marker is not derived from a public hard-coded key.
    """
    import hashlib
    import hmac
    if not secret_seed:
        raise ValueError("a leak-marker secret is required")
    key = secret_seed.encode("utf-8")
    msg = recipient.strip().lower().encode("utf-8")
    digest = hmac.new(key, msg, hashlib.sha256).hexdigest()[:32]
    binary = "".join(f"{int(c, 16):04b}" for c in digest)
    chars = ["\uFEFF"]
    for bit in binary:
        chars.append("\u200C" if bit == "1" else "\u200B")
    chars.append("\uFEFF")
    return "".join(chars), digest


def _decode_zero_width_tag(text: str) -> str | None:
    """Extract and decode a zero-width hex signature from text string."""
    if not text or "\uFEFF" not in text:
        return None
    m = re.search(r"\uFEFF([\u200B\u200C]{128})\uFEFF", text)
    if not m:
        return None
    binary = m.group(1)
    bits = [1 if c == "\u200C" else 0 for c in binary]
    hex_chars = []
    for i in range(0, len(bits), 4):
        chunk = bits[i:i+4]
        val = sum(b << (3 - j) for j, b in enumerate(chunk))
        hex_chars.append(f"{val:x}")
    return "".join(hex_chars)


def _is_luhn_valid(number_str: str) -> bool:
    """Verify number_str with Luhn algorithm mod-10."""
    digits = [int(c) for c in number_str if c.isdigit()]
    if len(digits) < 13 or len(digits) > 19 or sum(digits) == 0:
        return False
    if digits[0] not in (2, 3, 4, 5, 6):
        return False
    checksum = 0
    reverse_digits = digits[::-1]
    for idx, d in enumerate(reverse_digits):
        if idx % 2 == 1:
            d *= 2
            if d > 9:
                d -= 9
        checksum += d
    return checksum % 10 == 0


def _is_iban_valid(iban_str: str) -> bool:
    """Verify IBAN according to ISO 7064 Modulo 97-10."""
    clean = re.sub(r"\s+", "", iban_str).upper()
    if len(clean) < 15 or len(clean) > 34 or not clean[:2].isalpha() or not clean[2:4].isdigit():
        return False
    reordered = clean[4:] + clean[:4]
    num_str = "".join(str(ord(c) - 55) if c.isalpha() else c for c in reordered)
    try:
        return int(num_str) % 97 == 1
    except ValueError:
        return False


def _is_valid_ssn(ssn_str: str) -> bool:
    """Verify US SSN does not have invalid area (000, 666, 900-999) or zero groups."""
    digits = re.sub(r"\D", "", ssn_str)
    if len(digits) != 9:
        return False
    area = int(digits[:3])
    group = int(digits[3:5])
    serial = int(digits[5:])
    if area in (0, 666) or (900 <= area <= 999) or group == 0 or serial == 0:
        return False
    return True


def _dispatch_protect_pdf(work: Path, source: Path, password: str,
                          owner_password: str, suffix: str = "",
                          open_password: str = "",
                          sign_compatible: bool = False,
                          stego_str: str = "",
                          watermark_text: str = "") -> Path:
    src = open_pdf(source, open_password)
    if stego_str and src.page_count > 0:
        # Invisible text mode 3 (neither fill nor stroke) on page 1
        src[0].insert_text(fitz.Point(10, 10), stego_str, render_mode=3)
    if watermark_text and src.page_count > 0:
        for page in src:
            rect = page.rect
            pt = fitz.Point(rect.width * 0.15, rect.height * 0.5)
            page.insert_text(
                pt,
                watermark_text,
                fontsize=14,
                color=(0.5, 0.5, 0.5),
                morph=(pt, fitz.Matrix(-45)),
                overlay=True,
                render_mode=0,
            )
    tag = suffix or "protected"
    dest = work / f"{stem(source)}_{tag}.pdf"
    if not sign_compatible:
        perm = fitz.PDF_PERM_ACCESSIBILITY | fitz.PDF_PERM_PRINT
        src.save(dest, encryption=fitz.PDF_ENCRYPT_AES_256,
                 owner_pw=owner_password, user_pw=password,
                 permissions=perm, garbage=4, deflate=True, clean=True)
        src.close()
        return dest

    # PyMuPDF writes the trailer's /Encrypt dictionary inline. That is widely
    # accepted by viewers, but a strict incremental-signing writer correctly
    # requires /Encrypt to be an indirect object. Normalise the clean document
    # through pyHanko's AES-256 writer before adding the final signature.
    prepared = work / f".{stem(source)}_{tag}_prepared.pdf"
    src.save(prepared, encryption=fitz.PDF_ENCRYPT_NONE,
             garbage=4, deflate=True, clean=True)
    src.close()
    try:
        from pyhanko.pdf_utils.crypt.permissions import StandardPermissions
        from pyhanko.pdf_utils.reader import PdfFileReader
        from pyhanko.pdf_utils.writer import copy_into_new_writer
        with prepared.open("rb") as doc_in, dest.open("wb") as encrypted_out:
            writer = copy_into_new_writer(PdfFileReader(doc_in))
            writer.encrypt(
                owner_password, password,
                perms=(StandardPermissions.ALLOW_PRINTING |
                       StandardPermissions.ALLOW_ASSISTIVE_TECHNOLOGY |
                       StandardPermissions.TOLERATE_MISSING_PDF_MAC),
                pdf_mac=False,
            )
            writer.write(encrypted_out)
    except ImportError:
        raise ToolError("Sign & Dispatch requires pyhanko on the native server")
    except Exception as exc:
        raise ToolError(f"could not create sign-compatible AES-256 PDF: {exc}")
    return dest


def _dispatch_sign_encrypted_pdf(work: Path, source: Path, password: str,
                                 certificate: Path, cert_password: str,
                                 field_name: str) -> Path:
    """Append a signature to an already encrypted PDF and verify final bytes."""
    try:
        from pyhanko.pdf_utils.incremental_writer import IncrementalPdfFileWriter
        from pyhanko.pdf_utils.reader import PdfFileReader
        from pyhanko.sign import signers, validation
        from pyhanko_certvalidator import ValidationContext
    except ImportError:
        raise ToolError("Sign & Dispatch requires pyhanko on the native server")
    if not cert_password:
        raise ToolError("certificate password is required")
    try:
        signer = signers.SimpleSigner.load_pkcs12(
            str(certificate), passphrase=cert_password.encode("utf-8"))
    except Exception as exc:
        raise ToolError(f"could not load signing certificate: {exc}")

    dest = work / f"{stem(source)}_signed.pdf"
    try:
        with source.open("rb") as doc_in, dest.open("wb") as out_file:
            previous = PdfFileReader(doc_in)
            auth = previous.decrypt(password)
            if not auth:
                raise ToolError("could not authenticate encrypted PDF before signing")
            writer = IncrementalPdfFileWriter(doc_in, prev=previous)
            signers.sign_pdf(
                writer,
                signers.PdfSignatureMetadata(field_name=field_name),
                signer=signer,
                output=out_file,
            )
        with dest.open("rb") as signed_in:
            reader = PdfFileReader(signed_in)
            reader.decrypt(password)
            signatures = list(reader.embedded_signatures)
            if not signatures:
                raise ToolError("signed output contains no digital signature")
            # Structural/CMS validation must not depend on the host OS trust
            # store (which can be unavailable in containers). Trust-chain
            # policy remains the recipient viewer's responsibility.
            status = validation.validate_pdf_signature(
                signatures[-1], signer_validation_context=ValidationContext(
                    trust_roots=[], allow_fetching=False))
            if not status.intact or not status.valid:
                raise ToolError("final encrypted PDF failed signature validation")
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"Sign & Dispatch failed: {exc}")
    return dest


def email_secure(work: Path, inputs: list[Path], p: dict) -> Result:
    """Encrypt one or more PDFs, optionally sign last, then dispatch safely."""
    import datetime
    import json
    import secrets
    import string
    import time
    import smtp_manager

    pdf_inputs = [path for path in inputs if path.suffix.lower() == ".pdf"]
    certificates = [path for path in inputs if path.suffix.lower() in (".p12", ".pfx")]
    if not pdf_inputs:
        raise ToolError("Secure Email Dispatch needs at least one PDF")
    if len(pdf_inputs) > smtp_manager.MAX_ATTACHMENTS:
        raise ToolError(f"Secure Email Dispatch accepts at most {smtp_manager.MAX_ATTACHMENTS} PDFs")

    try:
        recipient = smtp_manager.validate_recipient_email(p.get("recipient_email") or "")
    except ValueError as exc:
        raise ToolError(str(exc))

    pw_mode = str(p.get("password_mode") or "random")
    pw = str(p.get("custom_password") or "").strip()
    if pw_mode == "random" or not pw:
        alphabet = string.ascii_letters + string.digits + "!@#$%^&*"
        pw = "".join(secrets.choice(alphabet) for _ in range(16))
    if len(pw) < 4:
        raise ToolError("password must be at least 4 characters")

    html_body = str(p.get("email_body_html") or "").strip() or None
    if html_body and re.search(r"\{\{\s*password\s*\}\}", html_body, re.I):
        raise ToolError("Email #1 templates cannot contain {{password}}")

    dispatch_pages = str(p.get("dispatch_pages") or "").strip()
    if dispatch_pages and len(pdf_inputs) != 1:
        raise ToolError("Split & Dispatch accepts exactly one master PDF")
    sources = list(pdf_inputs)
    if dispatch_pages:
        sources = [_dispatch_extract_pages(
            work, pdf_inputs[0], dispatch_pages,
            str(p.get("pdf_open_password") or ""))]

    sign_requested = p.get("sign_dispatch") in (True, "true", "1", 1)
    if sign_requested and len(certificates) != 1:
        raise ToolError("Sign & Dispatch requires one .p12 or .pfx certificate")

    # Correlation marker only: it can be stripped or copied and is not
    # forensic proof. Prefer a stable deployment secret; use a fresh
    # per-dispatch secret when the deployment did not configure one.
    marker_secret = (
        os.environ.get("LEAK_TRACER_SECRET")
        or os.environ.get("API_KEY")
        or secrets.token_urlsafe(32)
    )
    zero_width_str, stego_hash = _encode_zero_width_tag(recipient, marker_secret)
    watermark_msg = ""
    if p.get("recipient_watermark") in (True, "true", "1", 1):
        today = datetime.date.today().isoformat()
        watermark_msg = f"Confidential \u2022 Dispatched to {recipient} \u2022 {today} [ID: {stego_hash}]"

    protected_paths: list[Path] = []
    for index, source in enumerate(sources):
        owner_pw = secrets.token_urlsafe(24)
        encrypted = _dispatch_protect_pdf(
            work, source, pw, owner_pw,
            f"protected_{index + 1}" if len(sources) > 1 else "protected",
            "" if dispatch_pages else str(p.get("pdf_open_password") or ""),
            sign_compatible=sign_requested,
            stego_str=zero_width_str,
            watermark_text=watermark_msg)
        if sign_requested:
            encrypted = _dispatch_sign_encrypted_pdf(
                work, encrypted, pw, certificates[0],
                str(p.get("password_cert") or ""),
                f"SquishSignature{index + 1}")
        protected_paths.append(encrypted)

    smtp_config = {}
    if p.get("smtp_profile_json"):
        try:
            smtp_config = json.loads(p.get("smtp_profile_json"))
        except Exception:
            raise ToolError("saved SMTP profile is invalid")
    if not smtp_config.get("server") and p.get("smtp_server_profile_id"):
        import env_manager
        try:
            smtp_config = env_manager.get_profile_for_dispatch(
                int(p["smtp_server_profile_id"]))
        except (KeyError, ValueError) as exc:
            raise ToolError(f"SMTP profile not found: {exc}")
    if not smtp_config.get("server"):
        smtp_config = {
            "server": str(p.get("mail_server") or os.environ.get("MAIL_SERVER") or ""),
            "port": int(p.get("mail_port") or os.environ.get("MAIL_PORT") or 587),
            "username": str(p.get("mail_username") or os.environ.get("MAIL_USERNAME") or ""),
            "password": str(p.get("mail_password") or os.environ.get("MAIL_PASSWORD") or ""),
            "from_name": str(p.get("mail_from_name") or os.environ.get("MAIL_FROM_NAME") or ""),
            "security": str(p.get("mail_security") or "starttls"),
        }
    if not smtp_config.get("server") and not p.get("smtp_pool_json"):
        raise ToolError("SMTP server is not configured. Please supply SMTP settings or unlock your vault.")

    smtp_pool = None
    if p.get("smtp_pool_json"):
        try:
            smtp_pool = json.loads(p.get("smtp_pool_json"))
        except Exception:
            pass

    subject = str(p.get("email_subject") or "").strip() or None
    email2_subj = str(p.get("email2_subject") or "").strip() or None
    email2_body_tmpl = str(p.get("email2_body") or "").strip() or None
    try:
        raw_delay = p.get("delay_seconds")
        delay_sec = smtp_manager.clamp_dispatch_delay(0.1 if raw_delay in (None, "") else raw_delay)
    except ValueError as exc:
        raise ToolError(str(exc))
    thread_emails = p.get("thread_emails") in (True, "true", "1", 1)
    key_delivery_mode = str(p.get("key_delivery_mode") or "email").lower()

    try:
        send_result = smtp_manager.send_dual_secure_email(
            smtp=smtp_pool or smtp_config, recipient=recipient, pdf_path=protected_paths,
            password=pw, subject=subject, html_body=html_body,
            email2_subject=email2_subj, email2_body=email2_body_tmpl,
            delay_seconds=delay_sec, thread_emails=thread_emails,
            key_delivery_mode=key_delivery_mode,
            plain_text_only=p.get("email_plain_text_only") in (True, "true", "1", 1),
        )
    except smtp_manager.DeliveryUncertainError as exc:
        raise ToolError(f"Email delivery outcome is uncertain: {exc}")
    except Exception as exc:
        raise ToolError(f"Email delivery failed -- nothing was sent: {exc}")

    attachment_names = [path.name for path in protected_paths]
    base = stem(pdf_inputs[0])
    receipt_data = {
        "status": "partial_failure" if send_result.get("partial") else "success",
        "recipient": recipient,
        "pdf_file": attachment_names[0],
        "attachments": attachment_names,
        "password_included": False,
        "step1_sent": True,
        "step2_sent": bool(send_result.get("step2_sent")),
        "key_delivery_mode": key_delivery_mode,
        "oob_required": key_delivery_mode in {"oob", "dual"},
        "signed": sign_requested,
        "pages": dispatch_pages or None,
        "stego_tag": stego_hash,
        "relay_used": send_result.get("relay_used"),
        "retries": send_result.get("retries", 0),
        "timestamp": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
        "message": send_result.get("message"),
    }
    if send_result.get("partial"):
        receipt_data["error"] = send_result.get("error")
        receipt_data["resend"] = {
            "recipient": recipient,
            "pdf_filename": attachment_names[0], "subject": subject,
            "email2_subject": email2_subj, "email2_body": email2_body_tmpl,
        }
    if os.environ.get("INCLUDE_PASSWORD_IN_RECEIPT") == "1":
        receipt_data["password"] = pw
        receipt_data["password_included"] = True
        if "resend" in receipt_data:
            receipt_data["resend"]["password"] = pw
    receipt_file = work / f"{base}_delivery_receipt.json"
    receipt_file.write_text(json.dumps(receipt_data, indent=2), encoding="utf-8")
    return Result(receipt_file, "application/json", receipt_file.name)



def unlock(work: Path, inputs: list[Path], p: dict) -> Result:
    """Strip password encryption from a document whose password you know."""
    src = open_pdf(inputs[0], p.get("password", ""))
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_unlocked.pdf", shrink=False)
    src.close()
    return Result(dest, PDF, f"{base}_unlocked.pdf")


def redact(work: Path, inputs: list[Path], p: dict) -> Result:
    """Permanently destroys matched content -- pixels and underlying text.

    apply_redactions rewrites the content stream, so the text cannot be
    recovered by selecting under the box (the classic fake-redaction failure).
    """
    terms = [t for t in (p.get("terms") or "").splitlines() if t.strip()]
    if not terms:
        raise ToolError("supply at least one term to redact, one per line")
    src = open_pdf(inputs[0], p.get("password", ""))
    fill = _hex_rgb(p.get("color") or "#000000")
    hits = 0
    for pno in parse_pages(p.get("pages", ""), src.page_count):
        page = src[pno]
        for term in terms:
            for rect in page.search_for(term.strip(), quads=False):
                page.add_redact_annot(rect, fill=fill)
                hits += 1
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)
    if not hits:
        src.close()
        raise ToolError("none of those terms were found in the text layer")
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_redacted.pdf")
    src.close()
    return Result(dest, PDF, f"{base}_redacted.pdf")


def auto_redact(work: Path, inputs: list[Path], p: dict) -> Result:
    """Automatically find and redact sensitive patterns (Emails, SSNs, Phones, CCs, IBANs, Dictionary terms)."""
    src = open_pdf(inputs[0], p.get("password", ""))
    fill = _hex_rgb(p.get("color") or "#000000")
    
    # Custom dictionary terms from uploaded secondary file or textarea
    custom_terms: list[str] = []
    if len(inputs) > 1 and inputs[1].suffix.lower() in (".txt", ".csv"):
        try:
            content = inputs[1].read_text(encoding="utf-8", errors="replace")
            for line in content.splitlines():
                for term in line.split(","):
                    t = term.strip().strip('"').strip("'")
                    if t and t not in custom_terms:
                        custom_terms.append(t)
                        if len(custom_terms) >= 5000:
                            break
        except Exception:
            pass
    if p.get("terms"):
        for line in str(p.get("terms")).splitlines():
            t = line.strip()
            if t and t not in custom_terms:
                custom_terms.append(t)
                if len(custom_terms) >= 5000:
                    break

    has_pattern = bool(
        p.get("redact_email") or p.get("redact_ssn") or
        p.get("redact_phone") or p.get("redact_cc") or
        p.get("redact_iban") or custom_terms
    )
    if not has_pattern:
        src.close()
        raise ToolError("select at least one type of information or supply dictionary terms to redact")
        
    hits = 0
    pages = parse_pages(p.get("pages", ""), src.page_count)
    for pno in pages:
        page = src[pno]
        text = page.get_text("text")
        page_hits = 0
        
        # 1. Emails
        if p.get("redact_email"):
            for match in set(re.findall(r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+", text)):
                for rect in page.search_for(match, quads=False):
                    page.add_redact_annot(rect, fill=fill)
                    hits += 1
                    page_hits += 1
                    
        # 2. SSNs (with area code validity check)
        if p.get("redact_ssn"):
            for match in set(re.findall(r"\b\d{3}[-\s]?\d{2}[-\s]?\d{4}\b", text)):
                if _is_valid_ssn(match):
                    for rect in page.search_for(match, quads=False):
                        page.add_redact_annot(rect, fill=fill)
                        hits += 1
                        page_hits += 1
                        
        # 3. Phone numbers
        if p.get("redact_phone"):
            for match in set(re.findall(r"\b(?:\+?1[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b", text)):
                for rect in page.search_for(match, quads=False):
                    page.add_redact_annot(rect, fill=fill)
                    hits += 1
                    page_hits += 1
                    
        # 4. Credit Cards (Luhn algorithm validated)
        if p.get("redact_cc"):
            for match in set(re.findall(r"\b(?:\d{4}[ -]?\d{4}[ -]?\d{4}[ -]?\d{1,7}|\d{13,19})\b", text)):
                if _is_luhn_valid(match):
                    for rect in page.search_for(match, quads=False):
                        page.add_redact_annot(rect, fill=fill)
                        hits += 1
                        page_hits += 1
                        
        # 5. IBAN Bank Accounts (ISO Modulo 97-10 validated)
        if p.get("redact_iban"):
            for match in set(re.findall(r"\b[A-Za-z]{2}\d{2}(?:[ -]?[A-Za-z0-9]{4}){2,7}(?:[ -]?[A-Za-z0-9]{1,4})?\b", text)):
                if _is_iban_valid(match):
                    for rect in page.search_for(match, quads=False):
                        page.add_redact_annot(rect, fill=fill)
                        hits += 1
                        page_hits += 1

        # 6. Custom Dictionary Terms
        for term in custom_terms:
            for rect in page.search_for(term, quads=False):
                page.add_redact_annot(rect, fill=fill)
                hits += 1
                page_hits += 1

        if page_hits > 0:
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)
            
    if not hits:
        src.close()
        raise ToolError("no matching patterns or terms were found in the text layer")
        
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_auto_redacted.pdf")
    src.close()
    return Result(dest, PDF, f"{base}_auto_redacted.pdf")


def compare(work: Path, inputs: list[Path], p: dict) -> Result:
    """Word-level diff of the two documents' text layers, as Markdown."""
    import difflib
    if len(inputs) != 2:
        raise ToolError("compare needs exactly 2 files")
    texts = []
    for f in inputs:
        d = open_pdf(f, p.get("password", ""))
        texts.append("\n".join(pg.get_text("text") for pg in d).splitlines())
        d.close()
    diff = difflib.unified_diff(texts[0], texts[1],
                                fromfile=inputs[0].name, tofile=inputs[1].name,
                                lineterm="", n=2)
    body = "\n".join(diff) or "(the two text layers are identical)"
    dest = work / "comparison.md"
    dest.write_text(f"# Comparison\n\n```diff\n{body}\n```\n", encoding="utf-8")
    return Result(dest, "text/markdown", "comparison.md")


# ------------------------------------------------------------- extract ---

def extract_images(work: Path, inputs: list[Path], p: dict) -> Result:
    """Pull out embedded images at their ORIGINAL resolution.

    Deliberately different from pdf-to-jpg, which rasterises whole pages at a
    chosen DPI. Here we recover the image objects the PDF actually carries, so
    a 4000px photo comes back as 4000px rather than downsampled to the page.
    """
    src = open_pdf(inputs[0], p.get("password", ""))
    base = stem(inputs[0])
    min_px = max(0, int(p.get("min_size", 64)))
    seen: set[int] = set()
    files: list[Path] = []
    for pno in parse_pages(p.get("pages", ""), src.page_count):
        for info in src[pno].get_images(full=True):
            xref = info[0]
            if xref in seen:      # the same logo on 200 pages is one file
                continue
            seen.add(xref)
            try:
                img = src.extract_image(xref)
            except Exception as exc:
                log.warning("xref %s not extractable: %s", xref, exc)
                continue
            if img["width"] < min_px or img["height"] < min_px:
                continue          # skip rules, bullets and spacer pixels
            f = work / f"{base}_p{pno + 1}_{xref}.{img['ext']}"
            f.write_bytes(img["image"])
            files.append(f)
    src.close()
    if not files:
        raise ToolError("no embedded images found above the minimum size")
    if len(files) == 1:
        return Result(files[0], f"image/{files[0].suffix.lstrip('.')}", files[0].name)
    dest = zip_dir(files, work / f"{base}_images.zip")
    return Result(dest, ZIP, f"{base}_images.zip")


def extract_attachments(work: Path, inputs: list[Path], p: dict) -> Result:
    src = open_pdf(inputs[0], p.get("password", ""))
    files: list[Path] = []
    for i in range(src.embfile_count()):
        info = src.embfile_info(i)
        name = safe_component(info.get("filename") or f"attachment_{i}")
        f = work / name
        f.write_bytes(src.embfile_get(i))
        files.append(f)
    src.close()
    if not files:
        raise ToolError("this PDF has no file attachments")
    if len(files) == 1:
        return Result(files[0], "application/octet-stream", files[0].name)
    dest = zip_dir(files, work / "attachments.zip")
    return Result(dest, ZIP, "attachments.zip")


def extract_fonts(work: Path, inputs: list[Path], p: dict) -> Result:
    """Recover embedded font files."""
    src = open_pdf(inputs[0], p.get("password", ""))
    base = stem(inputs[0])
    files: list[Path] = []
    seen: set[str] = set()
    for pno in parse_pages(p.get("pages", ""), src.page_count):
        for font in src.get_page_fonts(pno):
            xref = font[0]
            if xref in seen:
                continue
            seen.add(xref)
            try:
                font_name, ext, _, font_bytes = src.extract_font(xref)
            except Exception as exc:
                log.warning("xref %s font not extractable: %s", xref, exc)
                continue
            if not font_bytes:
                continue
            name = safe_component(f"{font_name}.{ext}")
            f = work / f"{base}_p{pno + 1}_{name}"
            f.write_bytes(font_bytes)
            files.append(f)
    src.close()
    if not files:
        raise ToolError("no embedded fonts found")
    if len(files) == 1:
        return Result(files[0], "application/octet-stream", files[0].name)
    dest = zip_dir(files, work / f"{base}_fonts.zip")
    return Result(dest, ZIP, f"{base}_fonts.zip")


def safe_component(name: str) -> str:
    # os.path.basename does NOT split on backslashes on Linux, so a Windows
    # path (`..\..\system32`) would arrive whole and its `..` survive. Normalise
    # both separators first, then drop leading dots so no `..` can remain.
    name = str(name).replace("\\", "/")
    name = os.path.basename(name)
    name = re.sub(r"[^\w.\- ]+", "_", name).lstrip(".")
    return name[:100] or "file"


# ------------------------------------------------------------ metadata ---

META_KEYS = ("title", "author", "subject", "keywords", "creator", "producer")


def metadata(work: Path, inputs: list[Path], p: dict) -> Result:
    """Edit or strip document metadata.

    Stripping also clears the XML metadata stream, which is the part people
    forget: clearing the info dictionary alone leaves author and software
    fingerprints sitting in the XMP packet.
    """
    base = stem(inputs[0])
    dest = work / f"{base}_metadata.pdf"

    # NOTE: do NOT use `exiftool -all=` to strip a PDF. exiftool cannot truly
    # delete PDF metadata -- PDF is incremental, so it appends an update and the
    # original info/XMP stays recoverable in the file (exiftool itself warns of
    # this). PyMuPDF's save() rewrites the document with garbage collection, so
    # the cleared info dict and XMP packet are actually gone. That rewrite is
    # the authoritative strip; exiftool would only add a false sense of safety.
    src = open_pdf(inputs[0], p.get("password", ""))
    if p.get("strip"):
        src.set_metadata({})
        src.del_xml_metadata()
    else:
        meta = dict(src.metadata or {})
        for k in META_KEYS:
            if p.get(k) is not None:
                meta[k] = str(p.get(k))
        src.set_metadata(meta)
    save(src, dest)
    src.close()
    return Result(dest, PDF, f"{base}_metadata.pdf")


def sign_pdf(work: Path, inputs: list[Path], p: dict) -> Result:
    """Sign a PDF using a PKCS#12 (.p12/.pfx) certificate."""
    pdfs = [f for f in inputs if f.suffix.lower() == ".pdf"]
    p12s = [f for f in inputs if f.suffix.lower() in (".p12", ".pfx")]
    if not pdfs or not p12s:
        raise ToolError("sign-pdf requires both a PDF and a .p12 certificate file")
    
    pdf_file = pdfs[0]
    p12_file = p12s[0]
    password = p.get("password_cert", "").encode("utf8")
    if not password:
        raise ToolError("certificate password is required")
        
    try:
        from pyhanko.sign import signers
        from pyhanko.pdf_utils.incremental_writer import IncrementalPdfFileWriter
    except ImportError:
        raise ToolError("pyhanko is not installed")
        
    try:
        signer = signers.SimpleSigner.load_pkcs12(str(p12_file), passphrase=password)
    except Exception as exc:
        raise ToolError(f"could not load certificate: {exc}")

    base = stem(pdf_file)
    dest = work / f"{base}_signed.pdf"

    # pyhanko's high-level signers.sign_pdf(writer, meta, signer=, output=)
    # writes the signed document to `output` and returns the stream -- there is
    # no .write_to_stream() on the result. Use a `with` for both handles so the
    # output file is closed even when signing raises.
    try:
        with open(pdf_file, "rb") as doc_in, open(dest, "wb") as out_file:
            w = IncrementalPdfFileWriter(doc_in)
            signers.sign_pdf(
                w,
                signers.PdfSignatureMetadata(field_name="Signature1"),
                signer=signer,
                output=out_file,
            )
    except Exception as exc:
        raise ToolError(f"signing failed: {exc}")

    return Result(dest, PDF, f"{base}_signed.pdf")


def verify_signature(work: Path, inputs: list[Path], p: dict) -> Result:
    """Verify signatures on a PDF document."""
    try:
        from pyhanko.sign import validation
        from pyhanko.pdf_utils.reader import PdfFileReader
    except ImportError:
        raise ToolError("pyhanko is not installed")
        
    pdf_file = inputs[0]
    out = []
    try:
        with open(pdf_file, 'rb') as doc_in:
            r = PdfFileReader(doc_in)
            if not r.embedded_signatures:
                raise ToolError("this document has no digital signatures")
                
            for sig in r.embedded_signatures:
                try:
                    # No trust context is supplied, so this reports cryptographic
                    # integrity (intact/valid), not a trusted-chain verdict --
                    # SimpleCertificateValidator is not a real pyhanko class.
                    status = validation.validate_pdf_signature(sig)
                    out.append(f"Signature '{sig.field_name}': "
                               f"{'INTACT' if status.intact else 'MODIFIED'}")
                    out.append(f"  Cryptographically valid: {status.valid}")
                    out.append(f"  Covers whole document: "
                               f"{not status.modification_level or status.modification_level.name}")
                    try:
                        out.append(f"  Signer: {status.signing_cert.subject.human_friendly}")
                    except Exception:
                        pass
                except Exception as exc:
                    out.append(f"Signature '{sig.field_name}': ERROR ({exc})")
    except ToolError:
        raise
    except Exception as exc:
        raise ToolError(f"could not read signatures: {exc}")
        
    base = stem(pdf_file)
    dest = work / f"{base}_signatures.txt"
    dest.write_text("\n".join(out) + "\n", encoding="utf-8")
    return Result(dest, "text/plain", f"{base}_signatures.txt")


def flatten(work: Path, inputs: list[Path], p: dict) -> Result:
    """Bake annotations and form fields into the page content.

    After this, a filled form cannot be un-filled and a highlight cannot be
    dragged off. Note this is about *editability*, not confidentiality -- to
    actually destroy content, use Redact.
    """
    src = open_pdf(inputs[0], p.get("password", ""))
    try:
        src.bake(annots=True, widgets=True)
    except AttributeError:
        raise ToolError("flatten needs PyMuPDF 1.24.6 or newer")
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_flat.pdf")
    src.close()
    return Result(dest, PDF, f"{base}_flat.pdf")


# ------------------------------------------------------------- imposition ---

def n_up(work: Path, inputs: list[Path], p: dict) -> Result:
    """Place 2 or 4 source pages on each output sheet."""
    src = open_pdf(inputs[0], p.get("password", ""))
    per = int(p.get("per_sheet", 2))
    if per not in (2, 4):
        raise ToolError("pages per sheet must be 2 or 4")
    gap = max(0, min(72, float(p.get("gap", 8))))
    first = src[0].rect
    if per == 2:
        # Two portrait pages side by side make a landscape sheet.
        sheet_w, sheet_h, cols, rows = first.height, first.width, 2, 1
    else:
        sheet_w, sheet_h, cols, rows = first.width, first.height, 2, 2

    out = fitz.open()
    pages = parse_pages(p.get("pages", ""), src.page_count)
    for i in range(0, len(pages), per):
        sheet = out.new_page(width=sheet_w, height=sheet_h)
        cell_w = (sheet_w - gap * (cols + 1)) / cols
        cell_h = (sheet_h - gap * (rows + 1)) / rows
        for slot, pno in enumerate(pages[i:i + per]):
            cx, cy = slot % cols, slot // cols
            box = fitz.Rect(
                gap + cx * (cell_w + gap), gap + cy * (cell_h + gap),
                gap + cx * (cell_w + gap) + cell_w,
                gap + cy * (cell_h + gap) + cell_h,
            )
            # keep_proportion stops a portrait page being stretched to fill a
            # landscape cell.
            sheet.show_pdf_page(box, src, pno, keep_proportion=True)
    src.close()
    base = stem(inputs[0])
    dest = save(out, work / f"{base}_{per}up.pdf")
    out.close()
    return Result(dest, PDF, f"{base}_{per}up.pdf")


def split_bookmarks(work: Path, inputs: list[Path], p: dict) -> Result:
    """Split at outline entries -- no page numbers to look up."""
    src = open_pdf(inputs[0], p.get("password", ""))
    level = max(1, int(p.get("level", 1)))
    toc = [t for t in src.get_toc() if t[0] <= level]
    if not toc:
        src.close()
        raise ToolError(f"no bookmarks at level {level} or above")

    starts = sorted({max(0, t[2] - 1) for t in toc})
    titles = {max(0, t[2] - 1): t[1] for t in toc}
    base = stem(inputs[0])
    parts: list[Path] = []
    for idx, start in enumerate(starts):
        end = (starts[idx + 1] - 1) if idx + 1 < len(starts) else src.page_count - 1
        if end < start:
            continue
        out = fitz.open()
        out.insert_pdf(src, from_page=start, to_page=end)
        label = safe_component(titles.get(start, f"section{idx + 1}"))[:60]
        part = save(out, work / f"{idx + 1:02d}_{label}.pdf")
        out.close()
        parts.append(part)
    src.close()
    if len(parts) == 1:
        return Result(parts[0], PDF, parts[0].name)
    dest = zip_dir(parts, work / f"{base}_sections.zip")
    return Result(dest, ZIP, f"{base}_sections.zip")


def header_footer(work: Path, inputs: list[Path], p: dict) -> Result:
    """Arbitrary running text. Generalises the page-numbers tool."""
    import datetime
    src = open_pdf(inputs[0], p.get("password", ""))
    template = str(p.get("text") or "").strip()
    if not template:
        raise ToolError("text is required")
    size = max(6, min(72, int(p.get("size", 10))))
    color = _hex_rgb(p.get("color") or "#555555")
    fx, fy = POSITIONS.get(p.get("position", "top"), (0.5, 0.08))
    pages = parse_pages(p.get("pages", ""), src.page_count)
    today = datetime.date.today().isoformat()
    for seq, pno in enumerate(pages):
        page = src[pno]
        label = (template
                 .replace("{n}", str(pno + 1))
                 .replace("{total}", str(src.page_count))
                 .replace("{date}", today)
                 .replace("{filename}", stem(inputs[0])))
        w = fitz.get_text_length(label, fontname="helv", fontsize=size)
        page.insert_text(fitz.Point(page.rect.width * fx - w / 2,
                                    page.rect.height * fy),
                         label, fontsize=size, fontname="helv",
                         color=color, overlay=True)
    base = stem(inputs[0])
    dest = save(src, work / f"{base}_headed.pdf")
    src.close()
    return Result(dest, PDF, f"{base}_headed.pdf")


def grayscale(work: Path, inputs: list[Path], p: dict) -> Result:
    """Convert to grayscale via Ghostscript. Often a large size win too."""
    if not inputs[0].stat().st_size:
        raise ToolError(f"cannot read {inputs[0].name}: file is empty")
    base = stem(inputs[0])
    dest = work / f"{base}_gray.pdf"
    run([
        "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.7",
        "-sColorConversionStrategy=Gray",
        "-dProcessColorModel=/DeviceGray",
        "-dOverrideICC=true", "-dNOPAUSE", "-dQUIET", "-dBATCH",
        f"-sOutputFile={dest}", str(inputs[0]),
    ], cwd=work)
    return Result(dest, PDF, f"{base}_gray.pdf")


def grayscale_wasm(work: Path, inputs: list[Path], p: dict) -> Result:
    """Browser-safe grayscale conversion by rebuilding pages as gray images.

    Ghostscript preserves vector text on the server. The browser fallback is
    intentionally explicit about rasterising: it trades search/accessibility
    for a real grayscale PDF without pretending WebAssembly has Ghostscript.
    """
    src = open_pdf(inputs[0], p.get("password", ""))
    pages = list(range(src.page_count))
    dpi = clamp_dpi(pages, src, 150)
    out = fitz.open()
    try:
        for pno in pages:
            page = src[pno]
            pix = page.get_pixmap(dpi=dpi, colorspace=fitz.csGRAY, alpha=False)
            new = out.new_page(width=page.rect.width, height=page.rect.height)
            new.insert_image(new.rect, pixmap=pix)
        base = stem(inputs[0])
        dest = save(out, work / f"{base}_gray.pdf")
    finally:
        out.close()
        src.close()
    return Result(dest, PDF, f"{base}_gray.pdf")


def rasterise(work: Path, inputs: list[Path], p: dict) -> Result:
    """Rebuild every page as a flat image.

    Blunt but effective: it destroys embedded JavaScript, form fields,
    hyperlinks and any selectable text, which makes it a decent sanitiser for
    a PDF from an untrusted source. It also destroys accessibility and search,
    so it is a deliberate trade rather than a default.
    """
    src = open_pdf(inputs[0], p.get("password", ""))
    dpi = max(72, min(400, int(p.get("dpi", 150))))
    pages = parse_pages(p.get("pages", ""), src.page_count)
    dpi = clamp_dpi(pages, src, dpi)
    out = fitz.open()
    for pno in pages:
        page = src[pno]
        pix = page.get_pixmap(dpi=dpi)
        new = out.new_page(width=page.rect.width, height=page.rect.height)
        new.insert_image(new.rect, pixmap=pix)
    src.close()
    base = stem(inputs[0])
    dest = save(out, work / f"{base}_flattened.pdf")
    out.close()
    return Result(dest, PDF, f"{base}_flattened.pdf")


def _hex_rgb(value: str) -> tuple[float, float, float]:
    # Strict: an empty or None colour is a caller bug, not black. Callers that
    # want a default pass it explicitly (`p.get("color") or "#000000"`), so a
    # blank field falls back there while a truly missing value is caught here.
    if value is None or str(value).strip() == "":
        raise ToolError("colour is required")
    v = str(value).strip().lstrip("#")
    if len(v) == 3:
        v = "".join(c * 2 for c in v)
    if len(v) != 6 or not re.fullmatch(r"[0-9a-fA-F]{6}", v):
        raise ToolError(f"bad colour: {value!r}")
    return tuple(int(v[i:i + 2], 16) / 255 for i in (0, 2, 4))  # type: ignore


# -------------------------------------------------------------- registry ---

@dataclass
class Tool:
    key: str
    name: str
    group: str
    blurb: str
    fn: Callable[[Path, list[Path], dict], Result]
    accept: str = ".pdf"
    multi: bool = False       # accepts more than one file
    min_files: int = 1
    fields: list[dict[str, Any]] | None = None


def F(name: str, kind: str, label: str, **kw) -> dict:
    return {"name": name, "kind": kind, "label": label, **kw}


# `picker` drives the visual page grid in the UI. "select" = a set of pages to
# act on (blank still means all); the UI keeps this field and the grid in sync
# two ways. remove-pages and organize declare their own variants below.
PAGES = F("pages", "text", "Pages", placeholder="all, or 1-3,7,10-",
          help="Blank means every page.", picker="select")

TOOLS: list[Tool] = [
    # -- organize
    Tool("merge", "Merge PDF", "Organize",
         "Combine several PDFs into one, in the order you arrange them.",
         merge, multi=True, min_files=2, fields=[
             F("output_name", "text", "Output name", placeholder="merged",
               help="Auto-filled from your files; edit if you like. "
                    "The .pdf extension is added automatically."),
         ]),
    Tool("split", "Split PDF", "Organize",
         "Pull out a range, or burst the file into separate documents.",
         split, fields=[
             F("mode", "select", "Mode", options=[
                 ["ranges", "Extract selected pages into one PDF"],
                 ["every", "One PDF per page"],
                 ["chunks", "Fixed-size groups"]], default="ranges"),
             PAGES,
             F("size", "number", "Pages per group", default=2, min=1, max=500),
         ]),
    Tool("remove-pages", "Remove pages", "Organize",
         "Delete the pages you name and keep everything else.",
         remove_pages, fields=[
             F("pages", "text", "Pages to remove", placeholder="e.g. 2 or 4-6",
               help="The pages to delete.", picker="remove"),
         ]),
    Tool("organize", "Reorder pages", "Organize",
         "Rebuild the document in an explicit page order.",
         organize, fields=[
             F("pages", "text", "New order", placeholder="3,1,2,4-",
               help="Pages appear in exactly this order.", picker="order"),
         ]),
    Tool("rotate", "Rotate PDF", "Organize",
         "Turn pages a quarter, half or three-quarter turn.",
         rotate, fields=[
             F("angle", "select", "Rotation", options=[
                 ["90", "90 clockwise"], ["180", "180"],
                 ["270", "90 counter-clockwise"]], default="90"),
             PAGES,
         ]),
    Tool("split-bookmarks", "Split at bookmarks", "Organize",
         "Break a long document into sections using its own outline.",
         split_bookmarks, fields=[
             F("level", "select", "Split at", options=[
                 ["1", "Top-level bookmarks only"],
                 ["2", "Level 2 and above"],
                 ["3", "Level 3 and above"]], default="1"),
         ]),
    Tool("n-up", "N-up / booklet", "Organize",
         "Print two or four pages per sheet to save paper.",
         n_up, fields=[
             F("per_sheet", "select", "Pages per sheet",
               options=[["2", "2 up"], ["4", "4 up"]], default="2"),
             F("gap", "number", "Gap (pt)", default=8, min=0, max=72),
             PAGES,
         ]),
    # -- optimize
    Tool("compress", "Compress PDF", "Optimize",
         "Shrink the file by downsampling images and deduplicating resources.",
         compress, fields=[
             F("level", "select", "Compression", options=[
                 ["low", "Less compression, best quality"],
                 ["recommended", "Recommended"],
                 ["extreme", "Extreme, smallest file"]], default="recommended"),
         ]),
    Tool("repair", "Repair PDF", "Optimize",
         "Rebuild a damaged cross-reference table and recover what is readable.",
         repair),
    Tool("grayscale", "Grayscale PDF", "Optimize",
         "Strip colour for cheaper printing, usually a big size win too.",
         grayscale),
    Tool("ocr", "OCR PDF", "Optimize",
         "Add a searchable, selectable text layer to a scan.",
         ocr, fields=[
             F("lang", "select", "Language", options=[
                 ["eng", "English"], ["fra", "French"], ["deu", "German"],
                 ["spa", "Spanish"], ["por", "Portuguese"], ["ita", "Italian"]],
               default="eng"),
             F("deskew", "checkbox", "Straighten crooked scans"),
             F("force", "checkbox", "Re-OCR pages that already have text"),
         ]),
    # -- convert
    Tool("jpg-to-pdf", "Image to PDF", "Convert to PDF",
         "Turn JPG, PNG or WEBP images into a PDF, one image per page.",
         jpg_to_pdf, accept=".jpg,.jpeg,.png,.webp,.bmp,.tif,.tiff",
         multi=True, fields=[
             F("size", "select", "Page size", options=[
                 ["fit", "Fit to each image"], ["a4", "A4"],
                 ["letter", "US Letter"]], default="fit"),
             F("margin", "number", "Margin (pt)", default=0, min=0, max=200),
         ]),
    Tool("office-to-pdf", "Office to PDF", "Convert to PDF",
         "Word, Excel and PowerPoint files rendered to PDF by LibreOffice.",
         office_to_pdf, accept=",".join(sorted(OFFICE_EXT)), multi=True),
    Tool("md-to-pdf", "Markdown to PDF", "Convert to PDF",
         "Render Markdown to PDF with faithful tables, code and page breaks. "
         "Upload a .md file or paste text below.",
         md_to_pdf, accept=".md,.markdown,.mdown,.txt", min_files=0, fields=[
             F("md_text", "textarea", "Markdown text",
               placeholder="# Title\n\nPaste or type Markdown here. Used only "
                           "when no file is uploaded."),
             F("page_size", "select", "Page size",
               options=[["a4", "A4"], ["letter", "US Letter"]], default="a4"),
             F("margin", "number", "Margin (mm)", default=18, min=0, max=50),
             F("title", "text", "Document title",
               help="Defaults to the file name."),
             F("allow_remote", "checkbox", "Allow remote images",
               help="Off by default. Fetches https images referenced in the "
                    "document; private/loopback addresses stay blocked."),
         ]),
    Tool("pdf-to-jpg", "PDF to image", "Convert from PDF",
         "Render each page as a JPG or PNG at the resolution you choose.",
         pdf_to_jpg, fields=[
             F("format", "select", "Format",
               options=[["jpg", "JPG"], ["png", "PNG"]], default="jpg"),
             F("dpi", "number", "Resolution (DPI)", default=150, min=36, max=600),
             PAGES,
         ]),
    Tool("pdf-to-word", "PDF to Word", "Convert from PDF",
         "Rebuild paragraphs, tables and images as an editable .docx.",
         pdf_to_word),
    Tool("pdf-to-excel", "PDF to Excel", "Convert from PDF",
         "Extract detected tables into a worksheet each.",
         pdf_to_excel, fields=[PAGES]),
    Tool("pdf-to-powerpoint", "PDF to PowerPoint", "Convert from PDF",
         "One slide per page, rendered full bleed.",
         pdf_to_powerpoint, fields=[
             F("dpi", "number", "Resolution (DPI)", default=150, min=72, max=300),
             PAGES,
         ]),
    Tool("pdf-to-pdfa", "PDF to PDF/A", "Convert from PDF",
         "Convert to the PDF/A-2b archival standard.", pdf_to_pdfa),
    Tool("pdf-to-markdown", "PDF to Markdown", "Convert from PDF",
         "Extract the text layer as Markdown, page by page.",
         pdf_to_markdown, fields=[PAGES]),
    Tool("pdf-to-text", "PDF to Text", "Convert from PDF",
         "Extract the raw text layer as a .txt file.",
         pdf_to_text, fields=[PAGES]),
    Tool("extract-images", "Extract images", "Convert from PDF",
         "Recover embedded images at their original resolution.",
         extract_images, fields=[
             F("min_size", "number", "Ignore images smaller than (px)",
               default=64, min=0, max=4000,
               help="Filters out rules, bullets and spacer pixels."),
             PAGES,
         ]),
    Tool("extract-attachments", "Extract attachments", "Convert from PDF",
         "Pull out files embedded inside the PDF.",
         extract_attachments),
    Tool("extract-fonts", "Extract fonts", "Convert from PDF",
         "Recover embedded font files.",
         extract_fonts, fields=[PAGES]),
    # -- edit
    Tool("watermark", "Add watermark", "Edit",
         "Stamp text across the page, once or tiled.",
         watermark, fields=[
             F("text", "text", "Watermark text", placeholder="CONFIDENTIAL",
               required=True),
             F("mode", "select", "Layout", options=[
                 ["single", "Once per page"], ["tile", "Tiled"]],
               default="single"),
             F("position", "select", "Position",
               options=[[k, k.replace("-", " ")] for k in POSITIONS],
               default="center"),
             F("size", "number", "Font size", default=42, min=6, max=200),
             F("opacity", "number", "Opacity", default=0.25, min=0.05, max=1,
               step=0.05),
             F("angle", "number", "Angle", default=45, min=-180, max=180),
             F("color", "color", "Colour", default="#7c5cff"),
             PAGES,
         ]),
    Tool("page-numbers", "Add page numbers", "Edit",
         "Number the pages, with your own format and placement.",
         page_numbers, fields=[
             F("format", "text", "Format", default="{n}",
               help="{n} is the number, {total} the count. e.g. Page {n} of {total}"),
             F("position", "select", "Position",
               options=[[k, k.replace("-", " ")] for k in POSITIONS],
               default="bottom"),
             F("start", "number", "Start at", default=1, min=0, max=99999),
             F("size", "number", "Font size", default=11, min=6, max=72),
             F("color", "color", "Colour", default="#000000"),
             PAGES,
         ]),
    Tool("crop", "Crop PDF", "Edit",
         "Trim margins off every page, measured in points.",
         crop, fields=[
             F("top", "number", "Top (pt)", default=0, min=0, max=1000),
             F("bottom", "number", "Bottom (pt)", default=0, min=0, max=1000),
             F("left", "number", "Left (pt)", default=0, min=0, max=1000),
             F("right", "number", "Right (pt)", default=0, min=0, max=1000),
             PAGES,
         ]),
    Tool("header-footer", "Header / footer", "Edit",
         "Running text at the top or bottom, with date and filename tokens.",
         header_footer, fields=[
             F("text", "text", "Text", required=True,
               placeholder="Draft — {date} — page {n} of {total}",
               help="Tokens: {n} {total} {date} {filename}"),
             F("position", "select", "Position",
               options=[[k, k.replace("-", " ")] for k in POSITIONS],
               default="top"),
             F("size", "number", "Font size", default=10, min=6, max=72),
             F("color", "color", "Colour", default="#555555"),
             PAGES,
         ]),
    Tool("flatten", "Flatten PDF", "Edit",
         "Bake annotations and form fields in so they can no longer be edited.",
         flatten),
    Tool("metadata", "Edit metadata", "Edit",
         "Change the title and author, or strip identifying metadata entirely.",
         metadata, fields=[
             F("strip", "checkbox", "Strip everything (ignores the fields below)"),
             F("title", "text", "Title"),
             F("author", "text", "Author"),
             F("subject", "text", "Subject"),
             F("keywords", "text", "Keywords"),
         ]),
    # -- security
    Tool("rasterise", "Rasterise PDF", "Security",
         "Rebuild every page as a flat image, destroying scripts and form fields.",
         rasterise, fields=[
             F("dpi", "number", "Resolution (DPI)", default=150, min=72, max=400),
             PAGES,
         ]),
    Tool("protect", "Protect PDF", "Security",
         "Encrypt with AES-256 and set what readers are allowed to do.",
         protect, fields=[
             F("password_new", "password", "New password", required=True),
             F("allow_copy", "checkbox", "Allow copying text"),
             F("allow_modify", "checkbox", "Allow editing and annotation"),
         ]),
    Tool("email-secure", "Secure Email Dispatch", "Security",
         "Encrypt and dispatch PDFs with optional bursting, signing, a removable leak marker, and out-of-band key delivery.",
         email_secure, accept=".pdf,.p12,.pfx", multi=True, min_files=1, fields=[
             F("recipient_email", "text", "Recipient email address", required=True,
               placeholder="recipient@example.com", help="Destination inbox for the secure document and password."),
             F("email_subject", "text", "Email subject",
               placeholder="Confidential Document", help="Optional subject line."),
             F("password_mode", "select", "Password mode", options=[
                 ("random", "Auto-generate strong 16-character password"),
                 ("manual", "Specify custom password"),
             ], default="random"),
             F("custom_password", "password", "Custom password",
               placeholder="Enter password (if manual mode selected)", help="Only used if manual password mode is chosen."),
             F("recipient_watermark", "checkbox", "Stamp recipient leak watermark"),
         ]),
    Tool("unlock", "Unlock PDF", "Security",
         "Strip encryption from a PDF whose password you know.",
         unlock),
    Tool("redact", "Redact PDF", "Security",
         "Permanently erase matching text -- not a black box drawn on top.",
         redact, fields=[
             F("terms", "textarea", "Terms to redact", required=True,
               placeholder="one term per line"),
             F("color", "color", "Box colour", default="#000000"),
             PAGES,
         ]),
    Tool("auto-redact", "Auto-redact", "Security",
         "Best-effort pattern redaction for text layers. WARNING: scans, split spans, ligatures, and layout can cause misses; always verify the output.",
         auto_redact, accept=".pdf,.txt,.csv", multi=True, fields=[
             F("redact_email", "checkbox", "Redact Emails"),
             F("redact_ssn", "checkbox", "Redact SSNs"),
             F("redact_phone", "checkbox", "Redact Phone Numbers"),
             F("redact_cc", "checkbox", "Redact Credit Cards (Luhn validated)"),
             F("redact_iban", "checkbox", "Redact IBAN Bank Accounts (ISO Mod-97)"),
             F("terms", "textarea", "Custom dictionary terms", placeholder="one term per line"),
             F("color", "color", "Box colour", default="#000000"),
             PAGES,
         ]),
    Tool("sign-pdf", "Sign PDF", "Security",
         "Digitally sign a PDF using your own .p12 or .pfx certificate.",
         sign_pdf, multi=True, min_files=2, accept=".pdf,.p12,.pfx", fields=[
             F("password_cert", "password", "Certificate password", required=True),
         ]),
    Tool("verify-signature", "Verify Signature", "Security",
         "Check the validity of digital signatures on a PDF.",
         verify_signature),
    Tool("compare", "Compare PDFs", "Security",
         "Line-by-line diff of two documents' text.",
         compare, multi=True, min_files=2),
]

REGISTRY = {t.key: t for t in TOOLS}
